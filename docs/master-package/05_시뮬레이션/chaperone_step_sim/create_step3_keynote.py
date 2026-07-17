#!/usr/bin/env python3
"""Keynote-friendly PPTX: Chaperone STEP 3 simulation proof"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os

NAVY = RGBColor(0x0B, 0x1F, 0x33)
TEAL = RGBColor(0x0D, 0x94, 0x88)
TEAL_DK = RGBColor(0x0A, 0x6B, 0x62)
MINT = RGBColor(0x5E, 0xE0, 0xC5)
CORAL = RGBColor(0xE8, 0x5D, 0x4C)
AMBER = RGBColor(0xE8, 0x9B, 0x0C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x2A, 0x36)
GRAY = RGBColor(0x5A, 0x6A, 0x75)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xEE, 0xF4, 0xF6)
SOFT_TEAL = RGBColor(0xD4, 0xF0, 0xEB)
SOFT_CORAL = RGBColor(0xFD, 0xE8, 0xE5)
LINE = RGBColor(0xD0, 0xDC, 0xE0)
GREEN = RGBColor(0x1E, 0x8E, 0x5A)
SOFT_GREEN = RGBColor(0xE0, 0xF5, 0xE8)

W, H = Inches(13.333), Inches(7.5)
BASE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(BASE, "figures")

def set_run(run, size=14, bold=False, color=DARK, font="Apple SD Gothic Neo"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        el = rPr.find(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}")
        if el is None:
            el = etree.SubElement(rPr, f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}")
        el.set("typeface", font)

def add_shape(slide, left, top, width, height, fill=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, left, top, width, height)
    s.line.fill.background()
    if fill is not None:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    try:
        s.adjustments[0] = 0.06
    except Exception:
        pass
    return s

def add_rect(slide, left, top, width, height, fill=None):
    return add_shape(slide, left, top, width, height, fill=fill, shape=MSO_SHAPE.RECTANGLE)

def add_text(slide, left, top, width, height, text, size=14, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    try:
        tf._bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}.get(valign, "t"))
    except Exception:
        pass
    return box

def add_paras(slide, left, top, width, height, lines, size=12, color=DARK, spacing=5):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        run = p.add_run()
        run.text = line
        set_run(run, size=size, color=color)
    return box

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def footer(slide, page, total):
    add_rect(slide, 0, Inches(7.22), W, Inches(0.28), fill=NAVY)
    add_text(slide, Inches(0.35), Inches(7.24), Inches(10), Inches(0.24),
             "IDSX · Chaperone STEP 3 simulation · Native preferential stabilization · Non-therapeutic claim",
             size=9, color=MINT)
    add_text(slide, Inches(11.5), Inches(7.24), Inches(1.5), Inches(0.24),
             f"{page}/{total}", size=9, color=WHITE, align=PP_ALIGN.RIGHT)

def section_bar(slide, title, sub=""):
    add_rect(slide, 0, 0, W, Inches(0.95), fill=NAVY)
    add_rect(slide, 0, Inches(0.95), W, Inches(0.06), fill=TEAL)
    add_text(slide, Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.4), title, size=22, bold=True, color=WHITE)
    if sub:
        add_text(slide, Inches(0.4), Inches(0.55), Inches(12.5), Inches(0.3), sub, size=12, color=MINT)

def pic(slide, name, left, top, width, height=None):
    path = os.path.join(FIG, name)
    if os.path.isfile(path):
        if height:
            slide.shapes.add_picture(path, left, top, width=width, height=height)
        else:
            slide.shapes.add_picture(path, left, top, width=width)
        return True
    return False

def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    built = []

    def fin(s):
        built.append(s)
        return s

    # 1 Cover
    s = blank(prs)
    add_rect(s, 0, 0, W, H, fill=NAVY)
    add_rect(s, 0, Inches(5.8), W, Inches(1.7), fill=TEAL_DK)
    add_text(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.4),
             "IDSX  ·  SIMULATION PROOF  ·  ONE STEP ONLY", size=14, bold=True, color=MINT)
    add_text(s, Inches(0.6), Inches(1.8), Inches(12), Inches(1.6),
             "샤페론 적용 스텝은 다양하다.\n그중 STEP 3 — Native 선택적 안정화만\n시뮬레이션으로 입증한다.",
             size=28, bold=True, color=WHITE)
    add_text(s, Inches(0.6), Inches(3.8), Inches(12), Inches(1.2),
             "모델: IDS S152I 접힘 페널티 · U ⇌ N · N+C ⇌ NC (Kd_N ≪ Kd_U)\n"
             "결과: [PC]=10 μM에서 folded pool 0.58 → 0.92  (1.59× 농축)\n"
             "대조: 비선호 결합(Kd 동일)은 농축 실패  ·  치료·세포 구제 claim 없음",
             size=14, color=SOFT_TEAL)
    add_text(s, Inches(0.6), Inches(6.1), Inches(12), Inches(0.9),
             "5FQL 구조 앵커 + 질량작용 평형 + 표면 클램프 프록시\nKeynote deck · 2026-07 · Thilo Yoo / IDSX",
             size=13, color=WHITE)
    fin(s)

    # 2 Why multi-step
    s = blank(prs)
    section_bar(s, "1 · 샤페론 적용 스텝이 다양한 이유", "한 약이 한 번에 전 경로를 증명하지 않는다")
    steps = [
        ("S1", "번역·합성", "분자 공급"),
        ("S2", "접힘 핵형성", "조립 시작"),
        ("S3", "Native 선택 안정화", "본 시뮬 ★"),
        ("S4", "ERQC / ERAD↓", "품질 통과"),
        ("S5", "골지·M6P", "표지·출하"),
        ("S6", "리소좀 배달", "현장 도착"),
        ("S7", "pH 해리·촉매", "억제 창"),
        ("S8", "GAG 제거", "임상 표현"),
    ]
    for i, (code, title, note) in enumerate(steps):
        col, row = i % 4, i // 4
        left = Inches(0.4 + col * 3.2)
        top = Inches(1.3 + row * 2.7)
        fill = SOFT_TEAL if code == "S3" else CARD
        line = TEAL if code == "S3" else LINE
        add_shape(s, left, top, Inches(3.0), Inches(2.4), fill=fill)
        add_rect(s, left, top, Inches(3.0), Inches(0.5), fill=TEAL if code == "S3" else NAVY)
        add_text(s, left + Inches(0.15), top + Inches(0.1), Inches(2.7), Inches(0.35),
                 code + (" ★ SIM" if code == "S3" else ""), size=14, bold=True, color=WHITE)
        add_text(s, left + Inches(0.15), top + Inches(0.75), Inches(2.7), Inches(0.7),
                 title, size=16, bold=True, color=NAVY if code != "S3" else TEAL_DK)
        add_text(s, left + Inches(0.15), top + Inches(1.5), Inches(2.7), Inches(0.6),
                 note, size=13, color=GRAY)
    fin(s)

    # 3 Selected step
    s = blank(prs)
    section_bar(s, "2 · 선택한 스텝: STEP 3 Native preferential stabilization",
                "ER 조건에서 접힘 평형 U ⇌ N 을 질량작용으로 이동")
    add_shape(s, Inches(0.4), Inches(1.25), Inches(6.2), Inches(5.6), fill=CARD)
    add_text(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(0.4),
             "메커니즘 (질량작용)", size=16, bold=True, color=TEAL)
    add_paras(s, Inches(0.6), Inches(1.95), Inches(5.8), Inches(4.6), [
        "U  ⇌  N          K_fold = exp(−ΔG_fold / RT)",
        "N + C  ⇌  NC     Kd_N  (Native에 세게)",
        "U + C  ⇌  UC     Kd_U  (Unfolded에 약하게)",
        "",
        "조건: Kd_N ≪ Kd_U  (preferential native binding)",
        "판독: folded pool = f(N) + f(NC)",
        "",
        "변이: ΔG_fold = ΔG_WT + ΔΔG_mut",
        "  ΔG_WT = −3.0 kcal/mol",
        "  ΔΔG_mut = +2.8 kcal/mol (생물리 스케일)",
        "  T = 310 K (ER 유사)",
        "",
        "화공 비유: 반응기 평형 조성 제어 밸브 1개",
    ], size=13, spacing=4)
    pic(s, "fig0_mechanism_scheme.png", Inches(6.8), Inches(1.3), Inches(6.1))
    fin(s)

    # 4 Why S152I needs this step
    s = blank(prs)
    section_bar(s, "3 · 왜 S152I에 STEP 3이 필요한가 (5FQL 구조 앵커)",
                "촉매 파괴가 아니라 코어 네트워크·패킹 결함")
    add_shape(s, Inches(0.35), Inches(1.2), Inches(5.9), Inches(5.7), fill=CARD)
    add_text(s, Inches(0.55), Inches(1.35), Inches(5.5), Inches(0.35),
             "5FQL 실측 재계산", size=15, bold=True, color=TEAL)
    add_paras(s, Inches(0.55), Inches(1.85), Inches(5.5), Inches(4.8), [
        "• Ser152 OG ··· Asp148 OD2  ≈  2.53 Å",
        "• Ser152 OG ··· Thr118 OG1  ≈  2.85 Å",
        "• CB 4.5 Å 내 heavy atom  =  26 (조밀)",
        "• RSA 0% · 촉매거리 10.27 Å (안전권)",
        "",
        "S152I 모델 국소 항:",
        "• H-bond 네트워크 상실 (OG 소실)",
        "• 부피 확장 soft-clash 페널티 (조대)",
        "• → 접힘 평형이 U 쪽으로 기운다",
        "",
        "∴ 개입 논리 = 152에 직접 결합이 아니라",
        "  N 상태를 붙잡아 K_app를 회복 (STEP 3)",
    ], size=13, spacing=5)
    pic(s, "fig5_local_strain_components.png", Inches(6.4), Inches(1.25), Inches(6.5))
    fin(s)

    # 5 Pipeline figure
    s = blank(prs)
    section_bar(s, "4 · 다단 경로 중 시뮬 타깃 명시", "초록 = 본 연구에서 입증하는 유일한 스텝")
    pic(s, "fig1_steps_overview.png", Inches(0.4), Inches(1.4), Inches(12.5))
    add_text(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.1),
             "S4–S8(ERAD↓, 배달, 활성, GAG)은 이 덱의 입증 범위 밖이다. "
             "STEP 3 성공은 하류의 필요조건에 가깝고 충분조건이 아니다.",
             size=14, color=GRAY)
    fin(s)

    # 6 Main result dose-response
    s = blank(prs)
    section_bar(s, "5 · 주 결과: 용량–반응 — folded pool 농축",
                "Preferential PC만 농축 · 비선호 대조는 실패")
    pic(s, "fig2_dose_response_folded.png", Inches(0.3), Inches(1.15), Inches(8.0))
    add_shape(s, Inches(8.5), Inches(1.25), Inches(4.4), Inches(5.5), fill=SOFT_GREEN)
    add_text(s, Inches(8.7), Inches(1.45), Inches(4.0), Inches(0.4),
             "Proof numbers", size=16, bold=True, color=GREEN)
    add_paras(s, Inches(8.7), Inches(2.0), Inches(4.0), Inches(4.4), [
        "변이 · PC 없음",
        "  folded = 0.580",
        "",
        "변이 · [PC]=10 μM",
        "  Kd_N=1, Kd_U=50 μM",
        "  folded = 0.922",
        "",
        "농축 배율",
        "  1.59 ×",
        "",
        "비선호 대조",
        "  Kd_N = Kd_U → 농축 없음",
        "",
        "→ STEP 3 적용 입증",
        "  (선택적 결합이 핵심)",
    ], size=13, spacing=3)
    fin(s)

    # 7 Species
    s = blank(prs)
    section_bar(s, "6 · 종 분포: N+NC가 competent folded pool",
                "샤페론은 U를 직접 고치지 않고 N을 붙잡아 당긴다")
    pic(s, "fig3_species_bars.png", Inches(0.5), Inches(1.3), Inches(12.2))
    fin(s)

    # 8 Enrichment
    s = blank(prs)
    section_bar(s, "7 · 입증 지표: Folded enrichment vs no-PC",
                "농도와 함께 단조 증가 — 질량작용 예측과 일치")
    pic(s, "fig6_enrichment.png", Inches(1.5), Inches(1.2), Inches(10.2))
    fin(s)

    # 9 Surface clamp
    s = blank(prs)
    section_bar(s, "8 · 구조 프록시: 표면 패치 클램프 (바인더 경로)",
                "같은 STEP 3을 알로스테릭 구속으로 보조하는 시나리오")
    pic(s, "fig4_surface_clamp_rmsf_energy.png", Inches(0.35), Inches(1.15), Inches(8.0))
    add_shape(s, Inches(8.5), Inches(1.25), Inches(4.4), Inches(5.5), fill=SOFT_TEAL)
    add_text(s, Inches(8.7), Inches(1.45), Inches(4.0), Inches(0.4),
             "클램프 스캔", size=16, bold=True, color=TEAL_DK)
    add_paras(s, Inches(8.7), Inches(2.0), Inches(4.0), Inches(4.5), [
        "패치: P150, Y151, S154…",
        "(152 직접 결합 아님)",
        "",
        "k_clamp 0 → 8",
        "RMSF_152:",
        "  1.24 → 0.62 Å",
        "  (≈ 50% 감소)",
        "",
        "네트워크 전달 γ=0.45",
        "코어 모듈 요동 억제",
        "",
        "해석: 표면 바인더도",
        "STEP 3( N 안정화 )을",
        "구조적으로 구현 가능",
        "(가설 프록시 층)",
    ], size=13, spacing=3)
    fin(s)

    # 10 Methods box
    s = blank(prs)
    section_bar(s, "9 · 방법 요약 · 파라미터 · 재현", "과장을 막기 위한 정직한 Methods 슬라이드")
    cards = [
        ("집단 모델", "2-state + ligand\n질량수지 bisection\nT=310 K, RT 명시"),
        ("파라미터", "ΔG_WT=−3.0\nΔΔG_mut=+2.8\nKd_N=1, Kd_U=50 μM"),
        ("구조", "5FQL chain A\nH-bond 거리 재계산\nsoft clash probe"),
        ("대조", "Kd_N=Kd_U\n비선호 결합\n농축 실패 확인"),
        ("스케일 주의", "EvoEF2 +46은\n순위용 raw 값\n집단모델에 미사용"),
        ("재현", "run_step3_simulation.py\ndata/*.csv\nfigures/*.png"),
    ]
    for i, (t, b) in enumerate(cards):
        col, row = i % 3, i // 3
        left = Inches(0.4 + col * 4.25)
        top = Inches(1.25 + row * 2.85)
        add_shape(s, left, top, Inches(4.05), Inches(2.6), fill=CARD)
        add_rect(s, left, top, Inches(4.05), Inches(0.5), fill=NAVY)
        add_text(s, left + Inches(0.15), top + Inches(0.1), Inches(3.7), Inches(0.35), t, size=14, bold=True, color=MINT)
        add_text(s, left + Inches(0.2), top + Inches(0.7), Inches(3.6), Inches(1.7), b, size=14, color=DARK)
    fin(s)

    # 11 Non-claims
    s = blank(prs)
    section_bar(s, "10 · 입증한 것 / 입증하지 않은 것", "이 경계가 과학적 신뢰")
    add_shape(s, Inches(0.4), Inches(1.3), Inches(6.1), Inches(5.5), fill=SOFT_GREEN)
    add_text(s, Inches(0.65), Inches(1.5), Inches(5.6), Inches(0.4),
             "입증 (STEP 3)", size=18, bold=True, color=GREEN)
    add_paras(s, Inches(0.65), Inches(2.1), Inches(5.6), Inches(4.4), [
        "✓ Preferential native binding 모델에서",
        "  folded pool이 농도 의존적으로 증가",
        "✓ 10 μM에서 1.59× 농축 (본 파라미터)",
        "✓ 비선호 결합 대조는 농축 실패",
        "✓ 5FQL 상 S152 모듈이 안정화 대상이",
        "  될 구조적 이유 (H-bond·패킹)",
        "✓ 표면 클램프가 코어 RMSF를 낮출 수",
        "  있는 네트워크 프록시",
    ], size=14, spacing=6)
    add_shape(s, Inches(6.8), Inches(1.3), Inches(6.1), Inches(5.5), fill=SOFT_CORAL)
    add_text(s, Inches(7.05), Inches(1.5), Inches(5.6), Inches(0.4),
             "비입증 (범위 밖)", size=18, bold=True, color=CORAL)
    add_paras(s, Inches(7.05), Inches(2.1), Inches(5.6), Inches(4.4), [
        "✗ 세포 ERAD 감소 실측",
        "✗ 리소좀 배달·활성 회복",
        "✗ GAG 감소·임상 효능",
        "✗ 특정 화합물이 IDS 약",
        "✗ 암브록솔 IDS 적용",
        "✗ AF/EvoEF2 절대 ΔΔG = 진실",
        "✗ 전 스텝(S1–S8) 일괄 증명",
        "✗ 환자 투여 근거",
    ], size=14, spacing=6)
    fin(s)

    # 12 Conclusion
    s = blank(prs)
    section_bar(s, "11 · 결론 한 장", "다단 중 한 밸브가 물리적으로 작동함을 보였다")
    add_shape(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(2.2), fill=SOFT_TEAL)
    add_text(s, Inches(0.8), Inches(1.65), Inches(11.7), Inches(1.7),
             "샤페론 적용 스텝은 S1–S8로 다양하다.\n"
             "본 시뮬레이션은 그중 STEP 3(Native 선택적 안정화)에 대해,\n"
             "preferential binding 질량작용이 folded pool을 1.59× 농축함을 입증한다.",
             size=18, bold=True, color=NAVY)
    add_paras(s, Inches(0.8), Inches(3.9), Inches(11.7), Inches(2.8), [
        "• 구조(5FQL S152 모듈)는 “왜 이 스텝이 필요한지”를 설명하고,",
        "• 집단 모델은 “그 스텝에서 개입이 무엇을 바꾸는지”를 숫자로 보여 주며,",
        "• 음성 대조는 “아무 결합이나 되는 것이 아님”을 가른다.",
        "• 하류 스텝(ERAD·배달·촉매)은 별도 게이트 — IDSX 사전등록 규율과 동일.",
        "",
        "다음: residual activity 문헌 · wet ΔTm · (선택) 실제 리간드 도킹 — 여전히 STEP 3 보강 층.",
    ], size=14, spacing=5)
    fin(s)

    # 13 Appendix paths
    s = blank(prs)
    section_bar(s, "부록 · 파일 위치", "재현 패키지")
    add_shape(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.5), fill=CARD)
    add_paras(s, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.0), [
        "Desktop/IDSX_Master_Package_2026/chaperone_step_sim/",
        "  ├── 00_STEP3_시뮬레이션_방법과_결과.md",
        "  ├── create_step3_keynote.py",
        "  ├── run_step3_simulation.py  (재실행)",
        "  ├── figures/  fig0–fig6",
        "  └── data/     simulation_summary.json, dose_response_*.csv, clamp_scan.csv",
        "",
        "Keynote/PPTX 출력:",
        "  IDSX_Chaperone_STEP3_Simulation_Proof.key / .pptx",
        "",
        "관련: Master Package 11 성공기준 사전등록 · 12 모달리티 트리 · Ser152 기전 문서",
    ], size=15, spacing=6)
    fin(s)

    # footers
    total = len(built)
    for i, slide in enumerate(built):
        if i == 0:
            continue
        footer(slide, i + 1, total)

    out = os.path.join(BASE, "IDSX_Chaperone_STEP3_Simulation_Proof.pptx")
    prs.save(out)
    print("Saved", out, "slides", total)
    return out

if __name__ == "__main__":
    build()
