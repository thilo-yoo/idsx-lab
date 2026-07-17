#!/usr/bin/env python3
"""Keynote-friendly PPTX: Chaperone STEP 7 pH-window simulation"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os

NAVY = RGBColor(0x0B, 0x1F, 0x33)
TEAL = RGBColor(0x0D, 0x94, 0x88)
TEAL_DK = RGBColor(0x0A, 0x6B, 0x62)
MINT = RGBColor(0x5E, 0xE0, 0xC5)
CORAL = RGBColor(0xE8, 0x5D, 0x4C)
AMBER = RGBColor(0xE8, 0x9B, 0x0C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x2A, 0x36)
GRAY = RGBColor(0x5A, 0x6A, 0x75)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_TEAL = RGBColor(0xD4, 0xF0, 0xEB)
SOFT_CORAL = RGBColor(0xFD, 0xE8, 0xE5)
SOFT_AMBER = RGBColor(0xFE, 0xF3, 0xD6)
SOFT_GREEN = RGBColor(0xE0, 0xF5, 0xE8)
GREEN = RGBColor(0x1E, 0x8E, 0x5A)
LINE = RGBColor(0xD0, 0xDC, 0xE0)

W, H = Inches(13.333), Inches(7.5)
BASE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(BASE, "figures")

def set_run(run, size=14, bold=False, color=DARK, font="Apple SD Gothic Neo"):
    run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color; run.font.name = font
    rPr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        el = rPr.find(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}")
        if el is None:
            el = etree.SubElement(rPr, f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}")
        el.set("typeface", font)

def add_shape(slide, left, top, width, height, fill=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, left, top, width, height)
    s.line.fill.background()
    if fill is not None:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    try: s.adjustments[0] = 0.06
    except Exception: pass
    return s

def add_rect(slide, left, top, width, height, fill=None):
    return add_shape(slide, left, top, width, height, fill=fill, shape=MSO_SHAPE.RECTANGLE)

def add_text(slide, left, top, width, height, text, size=14, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    set_run(run, size=size, bold=bold, color=color)
    try:
        tf._bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}.get(valign, "t"))
    except Exception: pass
    return box

def add_paras(slide, left, top, width, height, lines, size=12, color=DARK, spacing=5):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        run = p.add_run(); run.text = line
        set_run(run, size=size, color=color)
    return box

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def footer(slide, page, total):
    add_rect(slide, 0, Inches(7.22), W, Inches(0.28), fill=NAVY)
    add_text(slide, Inches(0.35), Inches(7.24), Inches(10), Inches(0.24),
             "IDSX · Chaperone STEP 7 simulation · Lysosomal pH-window · ≠ STEP 3 deck",
             size=9, color=MINT)
    add_text(slide, Inches(11.5), Inches(7.24), Inches(1.5), Inches(0.24),
             f"{page}/{total}", size=9, color=WHITE, align=PP_ALIGN.RIGHT)

def section_bar(slide, title, sub=""):
    add_rect(slide, 0, 0, W, Inches(0.95), fill=NAVY)
    add_rect(slide, 0, Inches(0.95), W, Inches(0.06), fill=AMBER)
    add_text(slide, Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.4), title, size=20, bold=True, color=WHITE)
    if sub:
        add_text(slide, Inches(0.4), Inches(0.55), Inches(12.5), Inches(0.3), sub, size=12, color=RGBColor(0xFE, 0xF3, 0xC7))

def pic(slide, name, left, top, width, height=None):
    path = os.path.join(FIG, name)
    if os.path.isfile(path):
        if height: slide.shapes.add_picture(path, left, top, width=width, height=height)
        else: slide.shapes.add_picture(path, left, top, width=width)
        return True
    return False

def build():
    prs = Presentation()
    prs.slide_width = W; prs.slide_height = H
    built = []
    def fin(s):
        built.append(s); return s

    # Cover
    s = blank(prs)
    add_rect(s, 0, 0, W, H, fill=NAVY)
    add_rect(s, 0, Inches(5.85), W, Inches(1.65), fill=RGBColor(0x92, 0x40, 0x0E))
    add_text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.35),
             "IDSX  ·  SECOND STEP DECK  ·  NOT A REPEAT OF STEP 3", size=13, bold=True, color=AMBER)
    add_text(s, Inches(0.6), Inches(1.7), Inches(12), Inches(1.8),
             "또 다른 샤페론 스텝:\nSTEP 7 — 리소좀 pH 창\n해리 · 촉매 가용 시뮬레이션",
             size=28, bold=True, color=WHITE)
    add_text(s, Inches(0.6), Inches(3.9), Inches(12), Inches(1.3),
             "억제형 PC: ER(pH 7.4)에서는 세게 붙고, 리소좀(pH 4.8)에서는 풀려 효소가 일한다.\n"
             "입증 지표: Kd 162× 증가 · ER 점유 0.89 · 리소좀 v/v_max 0.94 (C=5 μM)\n"
             "대조: pH 무관 강결합은 리소좀 활성 ≈ 0  ·  치료 claim 없음",
             size=14, color=RGBColor(0xFE, 0xF3, 0xC7))
    add_text(s, Inches(0.6), Inches(6.15), Inches(12), Inches(0.9),
             "IDSX · Inhibitory PC therapeutic window model · 2026-07\n"
             "Prior deck = STEP 3 (folding equilibrium). This deck = STEP 7 only.",
             size=12, color=WHITE)
    fin(s)

    # Why different step
    s = blank(prs)
    section_bar(s, "1 · 왜 STEP 3 다음이 아니라 STEP 7인가", "같은 ‘샤페론’이라도 밸브가 다르다")
    add_shape(s, Inches(0.4), Inches(1.25), Inches(6.1), Inches(5.5), fill=SOFT_TEAL)
    add_text(s, Inches(0.6), Inches(1.45), Inches(5.7), Inches(0.4), "STEP 3 (이전 덱)", size=16, bold=True, color=TEAL_DK)
    add_paras(s, Inches(0.6), Inches(2.0), Inches(5.7), Inches(4.4), [
        "질문: 접힌 분자 분율을 올릴 수 있나?",
        "도구: U ⇌ N, Kd_N ≪ Kd_U",
        "출력: folded pool 1.59× 농축",
        "공정 위치: ER 접힘 반응기 평형",
        "",
        "한계: 리소좀에서 약이 안 떨어지면",
        "접힌 채로도 촉매 불가",
    ], size=14, spacing=6)
    add_shape(s, Inches(6.8), Inches(1.25), Inches(6.1), Inches(5.5), fill=SOFT_AMBER)
    add_text(s, Inches(7.0), Inches(1.45), Inches(5.7), Inches(0.4), "STEP 7 (본 덱)", size=16, bold=True, color=AMBER)
    add_paras(s, Inches(7.0), Inches(2.0), Inches(5.7), Inches(4.4), [
        "질문: 리소좀에서 촉매할 자유 효소가 남나?",
        "도구: pH-의존 Kd + MM 속도",
        "출력: ER 점유↑ & lys v↑ 동시",
        "공정 위치: 현장 반응기 억제 해제",
        "",
        "한계: 배달·재고가 0이면",
        "창이 열려도 처리량 0",
    ], size=14, spacing=6)
    fin(s)

    # Cascade
    s = blank(prs)
    section_bar(s, "2 · 다단 경로 — 본 덱은 주황 S7만", "민트 S3 = 이전 프레젠테이션")
    pic(s, "fig1_steps_S7.png", Inches(0.35), Inches(1.25), Inches(12.5))
    add_text(s, Inches(0.5), Inches(5.9), Inches(12.2), Inches(1.0),
             "S4–S6·S8은 여전히 미시뮬. STEP 7 시사 ≠ 전 경로 증명.",
             size=14, color=GRAY)
    fin(s)

    # Mechanism
    s = blank(prs)
    section_bar(s, "3 · STEP 7 메커니즘 스키마", "Inhibitory PC therapeutic window")
    pic(s, "fig0_pH_window_scheme.png", Inches(0.8), Inches(1.2), Inches(11.7))
    fin(s)

    # Kd vs pH
    s = blank(prs)
    section_bar(s, "4 · pH에 따른 Kd 이동 (모델)", "α=0.85 · Lys/ER Kd 비 ≈ 162×")
    pic(s, "fig2_Kd_vs_pH.png", Inches(0.4), Inches(1.15), Inches(7.8))
    add_shape(s, Inches(8.4), Inches(1.25), Inches(4.5), Inches(5.5), fill=SOFT_AMBER)
    add_text(s, Inches(8.6), Inches(1.45), Inches(4.1), Inches(0.4), "수식", size=15, bold=True, color=AMBER)
    add_paras(s, Inches(8.6), Inches(2.0), Inches(4.1), Inches(4.5), [
        "log10 Kd(pH) =",
        "  log10 Kd_ref",
        "  + α (pH_ref − pH)",
        "",
        "pH_ref = 7.4",
        "Kd_ref = 0.5 μM",
        "α = 0.85",
        "",
        "pH↓ → Kd↑",
        "(산성에서 약하게 붙음)",
        "",
        "현상론 모델",
        "원자 양성자화 경로 ≠",
    ], size=13, spacing=4)
    fin(s)

    # Occupancy vs activity
    s = blank(prs)
    section_bar(s, "5 · 주 결과: ER 점유 vs 리소좀 활성", "같은 농도 축에서 두 목표가 공존하는 구간")
    pic(s, "fig3_occupancy_vs_activity.png", Inches(0.3), Inches(1.15), Inches(8.1))
    add_shape(s, Inches(8.55), Inches(1.25), Inches(4.35), Inches(5.5), fill=SOFT_GREEN)
    add_text(s, Inches(8.75), Inches(1.45), Inches(4.0), Inches(0.4), "C = 5 μM", size=16, bold=True, color=GREEN)
    add_paras(s, Inches(8.75), Inches(2.0), Inches(4.0), Inches(4.5), [
        "pH-window PC",
        "  ER occupancy 0.89",
        "  Lys v/v_max   0.94",
        "",
        "pH-insensitive tight",
        "  Lys v/v_max ≈ 0.01",
        "",
        "→ STEP 7 적용 시사:",
        "  창이 있어야",
        "  촉매 가용이 산다",
        "",
        "빨간 점선 = 실패 모드",
        "  (과억제)",
    ], size=13, spacing=4)
    fin(s)

    # Window index
    s = blank(prs)
    section_bar(s, "6 · Therapeutic window index", "occ_ER × v_lys — 너무 적어도, 너무 많아도 실패")
    pic(s, "fig4_window_index.png", Inches(1.2), Inches(1.15), Inches(10.8))
    fin(s)

    # Compartment bars
    s = blank(prs)
    section_bar(s, "7 · 동일 농도, 다른 STEP 7 결과", "창 있는 PC vs 창 없는 강결합")
    pic(s, "fig5_compartment_bars.png", Inches(0.8), Inches(1.2), Inches(11.6))
    fin(s)

    # Heatmap
    s = blank(prs)
    section_bar(s, "8 · 촉매 속도 지형 (pH × 농도)", "흰색 가로선: ER 7.4 / Lys 4.8")
    pic(s, "fig8_rate_heatmap.png", Inches(1.5), Inches(1.15), Inches(10.2))
    fin(s)

    # Alpha sensitivity
    s = blank(prs)
    section_bar(s, "9 · pH 민감도 α 민감도 분석", "α→0 이면 window 붕괴 — STEP 7의 필요조건")
    pic(s, "fig6_alpha_sensitivity.png", Inches(1.5), Inches(1.2), Inches(10.2))
    fin(s)

    # Bridge
    s = blank(prs)
    section_bar(s, "10 · STEP 3 × STEP 7 결합 스케치", "재고 밸브 × 촉매 가용 밸브")
    pic(s, "fig7_step3_step7_bridge.png", Inches(0.4), Inches(1.15), Inches(7.6))
    add_shape(s, Inches(8.2), Inches(1.25), Inches(4.7), Inches(5.5), fill=CARD)
    add_text(s, Inches(8.4), Inches(1.45), Inches(4.3), Inches(0.4), "Throughput 개념", size=15, bold=True, color=NAVY)
    add_paras(s, Inches(8.4), Inches(2.0), Inches(4.3), Inches(4.5), [
        "≈ 재고 분율 × v_lys",
        "",
        "PC 없음: 0.58 × 1",
        "  → 0.58",
        "",
        "pH-window:",
        "  0.92 × 0.94 → 0.87",
        "",
        "강결합 무창:",
        "  0.92 × 0.01 → 0.01",
        "",
        "교훈: STEP 3만 또는",
        "STEP 7만으로는 부족할 수",
        "있다 (개념 스케치).",
    ], size=13, spacing=3)
    fin(s)

    # Methods
    s = blank(prs)
    section_bar(s, "11 · 방법 · 대조 · 한계", "정직한 Methods")
    items = [
        ("결합", "1:1 E+C⇌EC\n이차방정식 해\nEC는 촉매 0"),
        ("속도", "MM: v=kcat E_free\nS/(Km+S)\n상대 v/v_max"),
        ("pH–Kd", "log-linear α 모델\n현상론\n원자 경로 아님"),
        ("대조", "α=0 초강결합\n리소좀 활성 붕괴\nwindow 실패"),
        ("입증 범위", "STEP 7 시사\n세포·GAG 없음\nIDS 특화 리간드 아님"),
        ("파일", "data/*.csv\nfigures/fig0–8\nsummary.json"),
    ]
    for i, (t, b) in enumerate(items):
        col, row = i % 3, i // 3
        left = Inches(0.4 + col * 4.25)
        top = Inches(1.25 + row * 2.85)
        add_shape(s, left, top, Inches(4.05), Inches(2.6), fill=CARD)
        add_rect(s, left, top, Inches(4.05), Inches(0.5), fill=NAVY)
        add_text(s, left + Inches(0.15), top + Inches(0.1), Inches(3.7), Inches(0.35), t, size=14, bold=True, color=AMBER)
        add_text(s, left + Inches(0.2), top + Inches(0.7), Inches(3.6), Inches(1.7), b, size=14, color=DARK)
    fin(s)

    # Claim boundary
    s = blank(prs)
    section_bar(s, "12 · 시사한 것 / 아닌 것", "STEP 7 한정")
    add_shape(s, Inches(0.4), Inches(1.3), Inches(6.1), Inches(5.5), fill=SOFT_GREEN)
    add_text(s, Inches(0.65), Inches(1.5), Inches(5.6), Inches(0.4), "시사 (supported)", size=18, bold=True, color=GREEN)
    add_paras(s, Inches(0.65), Inches(2.15), Inches(5.6), Inches(4.4), [
        "✓ pH-민감 억제형 모델에서",
        "  ER 점유와 리소좀 활성이",
        "  같은 농도대에 공존 가능",
        "✓ Kd 산성 이동(~162×)이 창을 만듦",
        "✓ 무창 강결합은 촉매 가용 파괴",
        "✓ window index에 최적 농도 존재",
        "✓ α 민감도 = STEP 7 필요조건",
        "✓ STEP 3과 직교하는 밸브임을 스케치",
    ], size=14, spacing=5)
    add_shape(s, Inches(6.8), Inches(1.3), Inches(6.1), Inches(5.5), fill=SOFT_CORAL)
    add_text(s, Inches(7.05), Inches(1.5), Inches(5.6), Inches(0.4), "아님 (out of scope)", size=18, bold=True, color=CORAL)
    add_paras(s, Inches(7.05), Inches(2.15), Inches(5.6), Inches(4.4), [
        "✗ 실제 IDS 저분자 후보 확정",
        "✗ 세포 리소좀 활성 실측",
        "✗ GAG 감소·임상 효능",
        "✗ STEP 3 재입증",
        "✗ 배달·M6P 해결",
        "✗ 암브록솔 = IDS 약",
        "✗ 원자 수준 pKa 경로 규명",
        "✗ 전 경로(S1–S8) 일괄 증명",
    ], size=14, spacing=5)
    fin(s)

    # Conclusion
    s = blank(prs)
    section_bar(s, "13 · 결론", "또 다른 스텝, 또 다른 입증")
    add_shape(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(2.4), fill=SOFT_AMBER)
    add_text(s, Inches(0.8), Inches(1.65), Inches(11.7), Inches(1.9),
             "샤페론 적용 스텝은 다양하다.\n"
             "이전 덱은 STEP 3(접힘 평형)을, 본 덱은 STEP 7(pH 창 촉매 가용)을\n"
             "각각 시뮬레이션으로 시사한다 — 서로 대체되지 않는다.",
             size=18, bold=True, color=NAVY)
    add_paras(s, Inches(0.8), Inches(4.1), Inches(11.7), Inches(2.7), [
        "• Proof: C=5 μM에서 ER occupancy 0.89 & lysosomal v/v_max 0.94 (pH-window 모델).",
        "• 실패 대조: pH 무관 강결합 → 리소좀 활성 붕괴 (과억제).",
        "• 화공 메시지: 상류 재고 밸브(STEP 3)와 하류 억제 해제 밸브(STEP 7)를 분리 설계·검증할 것.",
        "• IDSX 다음: 실제 리간드 pH-DSF / 활성 assay로 STEP 7 가설을 wet으로 옮기기.",
    ], size=14, spacing=6)
    fin(s)

    total = len(built)
    for i, slide in enumerate(built):
        if i == 0: continue
        footer(slide, i + 1, total)

    out = os.path.join(BASE, "IDSX_Chaperone_STEP7_pHWindow_Proof.pptx")
    prs.save(out)
    print("Saved", out, "slides", total)
    return out

if __name__ == "__main__":
    build()
