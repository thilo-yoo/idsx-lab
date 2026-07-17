#!/usr/bin/env python3
"""Keynote deck: Full S1–S8 chaperone step simulations"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os, json

NAVY = RGBColor(0x0B, 0x1F, 0x33)
TEAL = RGBColor(0x0D, 0x94, 0x88)
MINT = RGBColor(0x5E, 0xE0, 0xC5)
CORAL = RGBColor(0xE8, 0x5D, 0x4C)
AMBER = RGBColor(0xE8, 0x9B, 0x0C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x2A, 0x36)
GRAY = RGBColor(0x5A, 0x6A, 0x75)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_TEAL = RGBColor(0xD4, 0xF0, 0xEB)
SOFT_CORAL = RGBColor(0xFD, 0xE8, 0xE5)
SOFT_AMBER = RGBColor(0xFE, 0xF3, 0xD6)
SOFT_GREEN = RGBColor(0xE0, 0xF5, 0xE8)
GREEN = RGBColor(0x1E, 0x8E, 0x5A)
PURPLE = RGBColor(0x6B, 0x4C, 0x9A)

W, H = Inches(13.333), Inches(7.5)
BASE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(BASE, "figures")
DATA = os.path.join(BASE, "data")

def set_run(run, size=14, bold=False, color=DARK, font="Apple SD Gothic Neo"):
    run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color; run.font.name = font
    rPr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        el = rPr.find(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}")
        if el is None:
            el = etree.SubElement(rPr, f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}")
        el.set("typeface", font)

def add_shape(slide, left, top, width, height, fill=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, left, top, width, height)
    s.line.fill.background()
    if fill is not None:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    try: s.adjustments[0] = 0.06
    except Exception: pass
    return s

def add_rect(slide, left, top, width, height, fill=None):
    return add_shape(slide, left, top, width, height, fill=fill, shape=MSO_SHAPE.RECTANGLE)

def add_text(slide, left, top, width, height, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return box

def add_paras(slide, left, top, width, height, lines, size=12, color=DARK, spacing=4):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        run = p.add_run(); run.text = line
        set_run(run, size=size, color=color)
    return box

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def footer(slide, page, total):
    add_rect(slide, 0, Inches(7.22), W, Inches(0.28), fill=NAVY)
    add_text(slide, Inches(0.35), Inches(7.24), Inches(10), Inches(0.24),
             "IDSX · Full cascade S1–S8 simulations · Conceptual / non-clinical", size=9, color=MINT)
    add_text(slide, Inches(11.5), Inches(7.24), Inches(1.5), Inches(0.24),
             f"{page}/{total}", size=9, color=WHITE, align=PP_ALIGN.RIGHT)

def section_bar(slide, title, sub=""):
    add_rect(slide, 0, 0, W, Inches(0.95), fill=NAVY)
    add_rect(slide, 0, Inches(0.95), W, Inches(0.06), fill=TEAL)
    add_text(slide, Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.4), title, size=20, bold=True, color=WHITE)
    if sub:
        add_text(slide, Inches(0.4), Inches(0.55), Inches(12.5), Inches(0.3), sub, size=12, color=MINT)

def pic(slide, name, left, top, width, height=None):
    path = os.path.join(FIG, name)
    if os.path.isfile(path):
        if height: slide.shapes.add_picture(path, left, top, width=width, height=height)
        else: slide.shapes.add_picture(path, left, top, width=width)

def build():
    with open(os.path.join(DATA, "all_steps_summary.json")) as f:
        master = json.load(f)
    steps = master["steps"]
    prs = Presentation()
    prs.slide_width = W; prs.slide_height = H
    built = []
    def fin(s):
        built.append(s); return s

    # Cover
    s = blank(prs)
    add_rect(s, 0, 0, W, H, fill=NAVY)
    add_rect(s, 0, Inches(5.9), W, Inches(1.6), fill=TEAL)
    add_text(s, Inches(0.55), Inches(1.2), Inches(12), Inches(0.35),
             "IDSX  ·  FULL CASCADE SIMULATION SUITE", size=14, bold=True, color=MINT)
    add_text(s, Inches(0.55), Inches(1.75), Inches(12), Inches(1.6),
             "샤페론 전 단계 S1–S8\n각 단계 구체 시뮬레이션",
             size=32, bold=True, color=WHITE)
    add_text(s, Inches(0.55), Inches(3.7), Inches(12), Inches(1.5),
             "합성 → 핵형성 → Native 안정화 → ERAD/수출 → M6P → 배달 → pH 창 촉매 → GAG 제거\n"
             "단계마다 모델·PC 역할·수치 지표·실패 모드를 분리 제시\n"
             "치료 claim 없음 · 의사결정·교육·우선순위용 공정 시뮬",
             size=14, color=SOFT_TEAL)
    add_text(s, Inches(0.55), Inches(6.15), Inches(12), Inches(0.9),
             "Companion decks: STEP3 · STEP7  ·  This deck: ALL steps + coupling\n2026-07 · IDSX Master Package",
             size=12, color=WHITE)
    fin(s)

    # Philosophy
    s = blank(prs)
    section_bar(s, "0 · 왜 전 단계를 각각 시뮬하는가", "한 스텝 성공 ≠ 치료 성공")
    add_shape(s, Inches(0.4), Inches(1.25), Inches(12.5), Inches(5.6), fill=CARD)
    add_paras(s, Inches(0.7), Inches(1.5), Inches(11.9), Inches(5.1), [
        "1. 샤페론 경로는 직렬 공정이다. 한 단의 효율 η_i 가 0이면 전체 처리량 Θ ≈ 0.",
        "2. 스텝마다 PC의 역할이 다르다 (S1: 무관 · S3: 평형 · S7: pH 창 · S5: 간접만).",
        "3. 같은 ‘결합’이라도 S3 성공 + S7 실패(과억제)면 종점 GAG는 오히려 악화될 수 있다.",
        "4. IDSX Hard EXCLUDE / CAUTION 규칙은 이 다단 맵의 어느 밸브가 막혔는지와 대응한다.",
        "5. 본 스위트는 각 단에 재현 가능한 최소 물리·동역학 모델을 붙여 ‘적용됨’을 수치로 시사한다.",
        "",
        "Non-claim: 세포 실험·임상 효능·특정 IDS 약물 구조 확정이 아니다.",
    ], size=15, spacing=8)
    fin(s)

    # Overview figure
    s = blank(prs)
    section_bar(s, "1 · 캐스케이드 한눈 지도", "S1–S8 각각 전용 모델")
    pic(s, "fig00_cascade_overview.png", Inches(0.3), Inches(1.2), Inches(12.6))
    fin(s)

    # Metrics dashboard
    s = blank(prs)
    section_bar(s, "2 · 전 단계 proof metric 대시보드", "각 패널 = 해당 스텝 시뮬의 핵심 숫자")
    pic(s, "fig01_metrics_dashboard.png", Inches(0.25), Inches(1.1), Inches(12.7))
    fin(s)

    # Individual steps S1-S8: model card + figure
    step_figs = [
        ("figS1_synthesis.png", "S1 합성 공급", "모델: dP/dt = r − kP", "PC 역할: 없음", "null NMD 천장 ~15% → EXCLUDE"),
        ("figS2_nucleation.png", "S2 접힘 핵형성", "모델: 생산적 vs 트랩 TST 속도", "PC: 장벽 ΔG‡ 감소", "f_prod 0.08→0.92 (−3 kcal)"),
        ("figS3_folded.png", "S3 Native 안정화", "모델: U⇌N + 선호 결합", "PC: Kd_N≪Kd_U", "folded 0.58→0.92 (1.59×)"),
        ("figS4_export.png", "S4 ERAD/수출", "모델: η = k_exp/(k_exp+k_ERAD)", "PC: f_fold↑ ± ERAD 차폐", "η 0.33→0.78 (shield)"),
        ("figS5_M6P.png", "S5 M6P 표지", "모델: η = ηmax f/(K+f)", "PC: 간접(화물 품질)", "glyco-dead 고정 0.05"),
        ("figS6_delivery.png", "S6 리소좀 배달", "모델: G→E→L ODE + leak", "PC: k_leak 감소", "도착 0.44→0.73"),
        ("figS7_pHwindow.png", "S7 pH 창 촉매", "모델: Kd(pH)+MM", "PC: 중성 점유/산성 해리", "occ 0.89 & v 0.94 vs tight 0.01"),
        ("figS8_GAG.png", "S8 GAG 제거", "모델: 기질 반응기 수지", "PC: 상류 결과의 종점", "GAG 38→33 (S3+S7)"),
    ]

    for i in range(0, 8, 2):
        s = blank(prs)
        a, b = step_figs[i], step_figs[i+1]
        section_bar(s, f"{3+i//2} · {a[1]}  &  {b[1]}", "좌·우 각 스텝 시뮬 결과")
        # left
        pic(s, a[0], Inches(0.25), Inches(1.15), Inches(6.3))
        add_text(s, Inches(0.35), Inches(5.55), Inches(6.1), Inches(1.4),
                 f"{a[2]}\n{a[3]}\n→ {a[4]}", size=12, color=DARK)
        # right
        pic(s, b[0], Inches(6.7), Inches(1.15), Inches(6.3))
        add_text(s, Inches(6.8), Inches(5.55), Inches(6.1), Inches(1.4),
                 f"{b[2]}\n{b[3]}\n→ {b[4]}", size=12, color=DARK)
        fin(s)

    # Detailed table slide
    s = blank(prs)
    section_bar(s, "7 · 스텝 카드 요약표", "모델 · PC 역할 · 핵심 수치")
    headers = ["Step", "모델 한 줄", "PC 역할", "핵심 결과"]
    rows = [
        ["S1", "합성 수지 r/k", "없음", "null 천장 15%"],
        ["S2", "핵형성 분배 TST", "장벽 감소", "f 0.08→0.92"],
        ["S3", "U⇌N 선호 결합", "N 선택 안정화", "1.59× folded"],
        ["S4", "export vs ERAD", "f↑/차폐", "η 0.33→0.78"],
        ["S5", "M6P MM형", "간접", "glyco-dead 불가"],
        ["S6", "수송 ODE", "leak↓", "0.44→0.73"],
        ["S7", "pH–Kd + MM", "억제 창", "v 0.94 vs 0.01"],
        ["S8", "GAG 반응기", "종점", "축적 완화/악화 갈림"],
    ]
    # header
    y0 = Inches(1.2)
    ws = [Inches(1.2), Inches(3.5), Inches(3.5), Inches(3.8)]
    xs = [Inches(0.4)]
    for w in ws[:-1]:
        xs.append(xs[-1] + w)
    for j, h in enumerate(headers):
        add_rect(s, xs[j], y0, ws[j], Inches(0.4), fill=NAVY)
        add_text(s, xs[j] + Inches(0.08), y0 + Inches(0.08), ws[j] - Inches(0.1), Inches(0.3), h, size=12, bold=True, color=MINT)
    for i, row in enumerate(rows):
        y = y0 + Inches(0.42) + i * Inches(0.65)
        bg = SOFT_TEAL if i % 2 == 0 else CARD
        for j, cell in enumerate(row):
            add_rect(s, xs[j], y, ws[j], Inches(0.6), fill=bg)
            add_text(s, xs[j] + Inches(0.08), y + Inches(0.12), ws[j] - Inches(0.12), Inches(0.4), cell, size=12, bold=(j == 0), color=DARK)
    fin(s)

    # Coupling
    s = blank(prs)
    section_bar(s, "8 · 전 단계 결합 처리량 Θ", "직렬 효율 곱 — 시나리오 비교")
    pic(s, "fig99_cascade_throughput.png", Inches(0.4), Inches(1.15), Inches(8.0))
    add_shape(s, Inches(8.55), Inches(1.25), Inches(4.4), Inches(5.5), fill=SOFT_AMBER)
    add_text(s, Inches(8.75), Inches(1.45), Inches(4.0), Inches(0.4), "읽는 법", size=15, bold=True, color=AMBER)
    add_paras(s, Inches(8.75), Inches(2.0), Inches(4.0), Inches(4.5), [
        "Θ ∝ Π η_i",
        "",
        "null → 0",
        "disease → 매우 낮음",
        "S3 only → 부분 회복",
        "S3+S7 → 추가 회복",
        "full stack → 더 개선",
        "WT → 상한",
        "",
        "한 단 실패 = 전체 실패",
        "과억제 S7 = 종점 악화",
        "(S8 시뮬 참고)",
    ], size=13, spacing=4)
    fin(s)

    # What PC can/cannot
    s = blank(prs)
    section_bar(s, "9 · 스텝별 할 수 있는 것 / 없는 것", "설계·자문용 체크")
    cans = [
        ("S1", "—", "null 공급 복구"),
        ("S2", "핵형성 경로 편향", "서열 교정"),
        ("S3", "folded 분율↑", "배달 보장"),
        ("S4", "수출 η↑", "ERAD 전면 차단 주장"),
        ("S5", "간접 M6P↑", "글리코 사이트 null"),
        ("S6", "수송 안정화", "BBB/CNS"),
        ("S7", "리소좀 촉매 가용", "재고 0에서 기적"),
        ("S8", "축적 완화(종점)", "단독 기전 증명"),
    ]
    add_rect(s, Inches(0.4), Inches(1.2), Inches(2.0), Inches(0.4), fill=NAVY)
    add_rect(s, Inches(2.4), Inches(1.2), Inches(5.2), Inches(0.4), fill=NAVY)
    add_rect(s, Inches(7.6), Inches(1.2), Inches(5.2), Inches(0.4), fill=NAVY)
    add_text(s, Inches(0.5), Inches(1.28), Inches(1.8), Inches(0.3), "Step", size=12, bold=True, color=MINT)
    add_text(s, Inches(2.5), Inches(1.28), Inches(5.0), Inches(0.3), "할 수 있음 (모델상)", size=12, bold=True, color=MINT)
    add_text(s, Inches(7.7), Inches(1.28), Inches(5.0), Inches(0.3), "못 함 / 금지 과대해석", size=12, bold=True, color=MINT)
    for i, (st, can, cannot) in enumerate(cans):
        y = Inches(1.65) + i * Inches(0.62)
        add_rect(s, Inches(0.4), y, Inches(2.0), Inches(0.55), fill=SOFT_TEAL if i%2==0 else CARD)
        add_rect(s, Inches(2.4), y, Inches(5.2), Inches(0.55), fill=SOFT_GREEN if i%2==0 else CARD)
        add_rect(s, Inches(7.6), y, Inches(5.2), Inches(0.55), fill=SOFT_CORAL if i%2==0 else CARD)
        add_text(s, Inches(0.5), y+Inches(0.12), Inches(1.8), Inches(0.35), st, size=13, bold=True, color=TEAL)
        add_text(s, Inches(2.55), y+Inches(0.12), Inches(5.0), Inches(0.35), can, size=12, color=DARK)
        add_text(s, Inches(7.75), y+Inches(0.12), Inches(5.0), Inches(0.35), cannot, size=12, color=DARK)
    fin(s)

    # Link to IDSX triage
    s = blank(prs)
    section_bar(s, "10 · IDSX triage 규칙과의 연결", "시뮬이 규칙을 정당화하는 지점")
    cards = [
        ("Hard EXCLUDE null", "S1 공급 천장·S2–S8 전부 0", CORAL, SOFT_CORAL),
        ("EXCLUDE 촉매 kill", "S7–S8 촉매항 복구 불가", CORAL, SOFT_CORAL),
        ("CAUTION 포켓", "S7 과억제 위험↑", AMBER, SOFT_AMBER),
        ("TIER1 코어 모듈", "S2–S4 주 작용 구간", TEAL, SOFT_TEAL),
        ("에피토프≠변이자리", "S3 표면 클램프 논리", TEAL, SOFT_TEAL),
        ("Glyco 주의", "S5 구제 불가 분기", AMBER, SOFT_AMBER),
        ("pH 창 필수", "S7 없으면 S8 악화 가능", PURPLE, RGBColor(0xED, 0xE6, 0xF5)),
        ("다단 게이트", "전 단계 Θ 곱으로 go/no-go", NAVY, CARD),
    ]
    for i, (t, b, ac, bg) in enumerate(cards):
        col, row = i % 4, i // 4
        left = Inches(0.35 + col * 3.25)
        top = Inches(1.25 + row * 2.85)
        add_shape(s, left, top, Inches(3.1), Inches(2.6), fill=bg)
        add_rect(s, left, top, Inches(3.1), Inches(0.55), fill=ac)
        add_text(s, left+Inches(0.12), top+Inches(0.12), Inches(2.85), Inches(0.4), t, size=13, bold=True, color=WHITE)
        add_text(s, left+Inches(0.15), top+Inches(0.8), Inches(2.8), Inches(1.5), b, size=14, color=DARK)
    fin(s)

    # Non-claim + conclusion
    s = blank(prs)
    section_bar(s, "11 · 결론 · Non-claim", "전 단계 맵이 주는 것")
    add_shape(s, Inches(0.4), Inches(1.25), Inches(12.5), Inches(2.3), fill=SOFT_TEAL)
    add_text(s, Inches(0.7), Inches(1.5), Inches(11.9), Inches(1.8),
             "S1부터 S8까지 각 단계에 구체 시뮬레이션을 붙였다.\n"
             "샤페론은 전 구간 만능이 아니라, 단마다 다른 밸브로 작동한다.\n"
             "IDSX는 그 밸브를 구조 triage로 고르는 의사결정 층이다.",
             size=17, bold=True, color=NAVY)
    add_paras(s, Inches(0.7), Inches(3.8), Inches(11.9), Inches(3.0), [
        "• 재현 패키지: chaperone_all_steps_sim/data/*.csv · figures/* · all_steps_summary.json",
        "• 심화 단독 덱: STEP3 (평형 농축) · STEP7 (pH 창) — 본 덱은 전 구간 통합",
        "• Non-claim: 임상 효능·IDS 약물 확정·원자 수준 절대 자유에너지 증명 아님",
        "• 다음 wet: residual activity · pH-DSF · 활성 washout 으로 S3/S7 파라미터 교정",
    ], size=14, spacing=7)
    fin(s)

    total = len(built)
    for i, slide in enumerate(built):
        if i == 0: continue
        footer(slide, i+1, total)

    out = os.path.join(BASE, "IDSX_Chaperone_AllSteps_S1-S8_Simulations.pptx")
    prs.save(out)
    print("Saved", out, "slides", total)
    return out

if __name__ == "__main__":
    build()
