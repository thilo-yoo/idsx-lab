#!/usr/bin/env python3
"""
IDSX Master Presentation — Advisor-grade + Personal research archive
16:9 Keynote-friendly · python-pptx · Apple SD Gothic Neo
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os
from copy import deepcopy

# ── Palette ──────────────────────────────────────────────────────────────
NAVY = RGBColor(0x0B, 0x1F, 0x33)
TEAL = RGBColor(0x0D, 0x94, 0x88)
TEAL_DK = RGBColor(0x0A, 0x6B, 0x62)
MINT = RGBColor(0x5E, 0xE0, 0xC5)
CORAL = RGBColor(0xE8, 0x5D, 0x4C)
AMBER = RGBColor(0xE8, 0x9B, 0x0C)
PURPLE = RGBColor(0x6B, 0x4C, 0x9A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF7, 0xF9, 0xFB)
LIGHT_BG = RGBColor(0xEE, 0xF4, 0xF6)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x5A, 0x6A, 0x75)
DARK = RGBColor(0x1A, 0x2A, 0x36)
SOFT_TEAL = RGBColor(0xD4, 0xF0, 0xEB)
SOFT_CORAL = RGBColor(0xFD, 0xE8, 0xE5)
SOFT_AMBER = RGBColor(0xFE, 0xF3, 0xD6)
SOFT_PURPLE = RGBColor(0xED, 0xE6, 0xF5)
SOFT_NAVY = RGBColor(0xE8, 0xEE, 0xF3)
LINE = RGBColor(0xD0, 0xDC, 0xE0)
GREEN = RGBColor(0x1E, 0x8E, 0x5A)
SOFT_GREEN = RGBColor(0xE0, 0xF5, 0xE8)

W, H = Inches(13.333), Inches(7.5)
OUT_DIR = os.path.dirname(os.path.abspath(__file__))
IMG = {
    "cascade": "/Users/bbc/Desktop/IDSX/Ser152Ile_Mechanism_Review/Fig1_Ser152Ile_mechanism_cascade.png",
    "rank": "/Users/bbc/Desktop/IDSX/Ser152Ile_Mechanism_Review/Fig2_Ser152Ile_plausibility_ranking.png",
    "pipeline": "/Users/bbc/Desktop/IDSX/Ser152Ile_Mechanism_Review/Fig3_patient_friendly_pipeline.png",
    "2axis": "/Users/bbc/Desktop/IDSX/Ser152Ile_Mechanism_Review/Fig4_hotspot_2axis_network_vs_catalytic.png",
    "af_overlay": "/Users/bbc/Desktop/IDSX/website/assets/img/AF3_structure_overlay_WT_S152I.png",
    "af_profiles": "/Users/bbc/Desktop/IDSX/website/assets/img/Fig_AF3_WT_vs_S152I_profiles.png",
    "af_zoom": "/Users/bbc/Desktop/IDSX/website/assets/img/Fig_AF3_zoom_module.png",
    "ctss": "/Users/bbc/Desktop/IDSX/ctss_scatter.png",
    "c422": "/Users/bbc/Desktop/IDSX/c422f_pocket.png",
}


def set_run(run, size=14, bold=False, color=DARK, font="Apple SD Gothic Neo"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        el = rPr.find(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}")
        if el is None:
            el = etree.SubElement(rPr, f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}")
        el.set("typeface", font)


def add_shape(slide, left, top, width, height, fill=None, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, left, top, width, height)
    s.line.fill.background()
    if fill is not None:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line is not None:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    try:
        s.adjustments[0] = 0.06
    except Exception:
        pass
    return s


def add_rect(slide, left, top, width, height, fill=None, line=None):
    return add_shape(slide, left, top, width, height, fill=fill, line=line, shape=MSO_SHAPE.RECTANGLE)


def add_text(slide, left, top, width, height, text, size=14, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font="Apple SD Gothic Neo"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    try:
        tf._bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}.get(valign, "t"))
    except Exception:
        pass
    return box


def add_paras(slide, left, top, width, height, lines, size=12, color=DARK, bold=False,
              spacing=6, align=PP_ALIGN.LEFT, font="Apple SD Gothic Neo"):
    """lines: list of str or (str, dict overrides)"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            text, opts = item
        else:
            text, opts = item, {}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", align)
        p.space_after = Pt(opts.get("spacing", spacing))
        p.space_before = Pt(0)
        run = p.add_run()
        run.text = text
        set_run(run, size=opts.get("size", size), bold=opts.get("bold", bold),
                color=opts.get("color", color), font=font)
    return box


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def footer(slide, page, total, section=""):
    add_rect(slide, Inches(0), Inches(7.22), W, Inches(0.28), fill=NAVY)
    label = f"IDSX · Structure-guided triage for IDS (P22304 / 5FQL)  ·  {section}" if section else \
            "IDSX · Structure-guided triage for IDS (P22304 / 5FQL)"
    add_text(slide, Inches(0.35), Inches(7.24), Inches(10), Inches(0.24),
             label, size=9, color=MINT, bold=False)
    add_text(slide, Inches(11.5), Inches(7.24), Inches(1.5), Inches(0.24),
             f"{page} / {total}", size=9, color=WHITE, align=PP_ALIGN.RIGHT)


def section_bar(slide, title, subtitle=""):
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.95), fill=NAVY)
    add_rect(slide, Inches(0), Inches(0.95), W, Inches(0.06), fill=TEAL)
    add_text(slide, Inches(0.4), Inches(0.18), Inches(12), Inches(0.4),
             title, size=22, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.4), Inches(0.55), Inches(12), Inches(0.3),
                 subtitle, size=12, color=MINT)


def card(slide, left, top, width, height, fill=CARD, line=LINE):
    return add_shape(slide, left, top, width, height, fill=fill, line=line)


def try_pic(slide, path, left, top, width, height=None):
    if path and os.path.isfile(path):
        if height:
            slide.shapes.add_picture(path, left, top, width=width, height=height)
        else:
            slide.shapes.add_picture(path, left, top, width=width)
        return True
    return False


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    slides_meta = []  # track for footer after count known
    built = []

    def fin(slide, section=""):
        built.append((slide, section))
        return slide

    # ══════════════════════════════════════════════════════════════════════
    # 0. COVER
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    add_rect(s, 0, 0, W, H, fill=NAVY)
    add_rect(s, 0, Inches(5.9), W, Inches(1.6), fill=TEAL_DK)
    add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.4),
             "IDSX  RESEARCH  ARCHIVE  ·  ADVISOR  BRIEFING", size=14, bold=True, color=MINT)
    add_text(s, Inches(0.6), Inches(1.85), Inches(12), Inches(1.2),
             "헌터 증후군 효소 IDS에 대한\n구조 기반 샤페론 적합성 triage와\n안정화 개입 논리의 개발 과정·현황·계획",
             size=28, bold=True, color=WHITE)
    add_text(s, Inches(0.6), Inches(4.4), Inches(12), Inches(0.8),
             "Iduronate-2-sulfatase (P22304) · PDB 5FQL · MPS II\n"
             "목표: catalytic competence 회복 가능성의 구조적 의사결정 층 구축\n"
             "비목표: 치료 성공 claim · 암브록솔 = IDS 약 · AF RMSD = 병원성 순위",
             size=13, color=SOFT_TEAL)
    add_text(s, Inches(0.6), Inches(6.15), Inches(8), Inches(0.9),
             "화학공학 전공 연구자 관점의 사고 궤적 공개 자료\n"
             "자문·기록·전달용 마스터 덱  ·  2026-07  ·  v2.0 Master",
             size=12, color=WHITE)
    add_text(s, Inches(9.5), Inches(6.3), Inches(3.3), Inches(0.7),
             "Professor / Expert\nConsultation Grade", size=12, bold=True, color=MINT, align=PP_ALIGN.RIGHT)
    fin(s, "Cover")

    # ══════════════════════════════════════════════════════════════════════
    # 1. AGENDA
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "목차 — 이 덱이 다루는 것", "요약이 아니라 ‘개발 과정·논리·상태·계획’을 한 패키지로 고정")
    items = [
        ("01", "왜 이 문제인가", "MPS II · 효소 공급 실패 · 화학공학적 문제 정의"),
        ("02", "사고의 궤적", "무엇을 버리고 무엇을 남겼는지 — 개발 연대기"),
        ("03", "로직 프레임", "네트워크 · competence · ER · 열안정=온도계"),
        ("04", "현재까지 구축물", "DB 797 · amenability · AF · ΔΔG · CTSS · 포털"),
        ("05", "대표 케이스", "Ser152 모듈 · Arg88 반례 · C422 트랙 긴장"),
        ("06", "유의미성·한계", "강한 층 / 열린 층 · non-claim 규율"),
        ("07", "앞으로의 계획", "90일 · 12개월 · 게이트 · 자문 질문"),
        ("08", "전달 문장", "누구에게나 말할 수 있는 고정 스크립트"),
    ]
    for i, (num, title, desc) in enumerate(items):
        col = i % 4
        row = i // 4
        left = Inches(0.4 + col * 3.2)
        top = Inches(1.35 + row * 2.7)
        card(s, left, top, Inches(3.0), Inches(2.4), fill=CARD, line=LINE)
        add_rect(s, left, top, Inches(0.12), Inches(2.4), fill=TEAL)
        add_text(s, left + Inches(0.3), top + Inches(0.25), Inches(2.5), Inches(0.4),
                 num, size=20, bold=True, color=TEAL)
        add_text(s, left + Inches(0.3), top + Inches(0.75), Inches(2.5), Inches(0.5),
                 title, size=16, bold=True, color=NAVY)
        add_text(s, left + Inches(0.3), top + Inches(1.35), Inches(2.5), Inches(0.85),
                 desc, size=11, color=GRAY)
    fin(s, "Agenda")

    # ══════════════════════════════════════════════════════════════════════
    # 2. ONE SENTENCE + CLAIM BOUNDARY
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "한 문장 정의 + Claim / Non-claim 경계", "이 경계를 어기면 전체 신뢰가 무너진다")
    card(s, Inches(0.4), Inches(1.25), Inches(12.5), Inches(1.5), fill=SOFT_TEAL, line=TEAL)
    add_text(s, Inches(0.65), Inches(1.4), Inches(12), Inches(1.2),
             "공개 변이 데이터와 IDS 결정구조(5FQL)를 결합해, 미스센스 변이를\n"
             "「샤페론·안정화 바인더로 접근할 가치가 있는 군」과 「구조적으로 배제·주의해야 할 군」으로 가르는\n"
             "재현 가능한 triage 지도(의사결정 층)를 만든다.  치료 성공 주장이 아니다.",
             size=15, bold=True, color=NAVY)

    # two columns
    card(s, Inches(0.4), Inches(3.0), Inches(6.1), Inches(3.9), fill=SOFT_GREEN, line=GREEN)
    add_text(s, Inches(0.65), Inches(3.15), Inches(5.6), Inches(0.4),
             "CLAIM 가능 (방어 가능)", size=15, bold=True, color=GREEN)
    add_paras(s, Inches(0.65), Inches(3.6), Inches(5.6), Inches(3.1), [
        "• 5FQL 실측 좌표 기반 RSA·촉매거리·이황화·Ca 주석을 전 변이셋에 일관 적용",
        "• 빈도 랭킹이 아닌 amenability proxy 규칙으로 재순위화",
        "• Hard EXCLUDE / CAUTION_POCKET / TIER1–3 분리",
        "• 변이 자리(core node) ≠ 결합 자리(surface epitope) 원칙",
        "• Ser152·Asp308을 구조 모듈 대표 케이스로 심화",
        "• AF 결과를 ‘좌표 유사 ≠ 안정·무해’ 규율로 해석",
        "• 암브록솔·미갈라스타트에서 성공 조건만 이식",
        "• 후속 PC·바인더 캠페인의 의사결정 층 제공",
    ], size=12, color=DARK, spacing=5)

    card(s, Inches(6.8), Inches(3.0), Inches(6.1), Inches(3.9), fill=SOFT_CORAL, line=CORAL)
    add_text(s, Inches(7.05), Inches(3.15), Inches(5.6), Inches(0.4),
             "NON-CLAIM (절대 말하지 않음)", size=15, bold=True, color=CORAL)
    add_paras(s, Inches(7.05), Inches(3.6), Inches(5.6), Inches(3.1), [
        "• 치료 바인더/약을 설계·검증했다",
        "• 잔존 활성·세포 트래픽킹을 실측했다",
        "• 암브록솔이 IDS에 직접 작용한다",
        "• AF RMSD / pLDDT만으로 병원성·구제 순위",
        "• docking 성공 = 기능 회복 (D2S0 교훈 반)",
        "• TIER1 = ‘약 된다’ 목록",
        "• 임상 처방용 amenable genotype 목록",
        "• 신경형 MPS II 전달 문제 해결",
    ], size=12, color=DARK, spacing=5)
    fin(s, "Boundary")

    # ══════════════════════════════════════════════════════════════════════
    # 3. WHY — PROCESS ENGINEERING VIEW
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "01 · 왜 이 문제인가 — 화학공학자의 문제 정의",
                "기질 축적이 아니라 ‘기능 효소 분자 유량’의 붕괴")
    # process flow
    stages = [
        ("전사·번역", "서열 설계도", SOFT_PURPLE, PURPLE),
        ("ER 접힘", "조립 라인", SOFT_TEAL, TEAL),
        ("ERQC", "품질관리", SOFT_AMBER, AMBER),
        ("골지·M6P", "표지·출하", SOFT_GREEN, GREEN),
        ("리소좀", "현장 반응기", SOFT_CORAL, CORAL),
    ]
    for i, (name, meta, bg, ac) in enumerate(stages):
        left = Inches(0.4 + i * 2.55)
        card(s, left, Inches(1.3), Inches(2.35), Inches(1.55), fill=bg, line=ac)
        add_text(s, left + Inches(0.15), Inches(1.5), Inches(2.05), Inches(0.45),
                 name, size=14, bold=True, color=ac, align=PP_ALIGN.CENTER)
        add_text(s, left + Inches(0.15), Inches(2.05), Inches(2.05), Inches(0.5),
                 meta, size=12, color=DARK, align=PP_ALIGN.CENTER)
        if i < 4:
            add_text(s, left + Inches(2.15), Inches(1.75), Inches(0.4), Inches(0.4),
                     "→", size=18, bold=True, color=NAVY)

    card(s, Inches(0.4), Inches(3.15), Inches(12.5), Inches(3.7), fill=CARD, line=LINE)
    add_text(s, Inches(0.65), Inches(3.3), Inches(12), Inches(0.35),
             "공정 사고(Process thinking)로 다시 쓰면", size=14, bold=True, color=NAVY)
    add_paras(s, Inches(0.65), Inches(3.75), Inches(12), Inches(2.9), [
        "• 최종 생산량(throughput) = 리소좀에서 작동하는 IDS 분자 수 × 분자당 활성  —  측정 활성 저하는 k_cat=0만이 아님",
        "• 미스센스 중 일부는 ‘반응기(활성부위) 파괴’가 아니라 ‘조립 수율↓ · QC 불합격 · 출하 실패’ 유형일 수 있음",
        "• 약리 샤페론(PC) = 조립 단계에서 접힘 평형을 유리하게 밀어 출하량을 올리는 공정 보조제 가설",
        "• 그러나 모든 결함이 같은 공정 처방으로 고쳐지지 않음 → 결함 모드(fault mode) 분류가 선행되어야 함",
        "• docking 성공 ≠ 기능 회복 (D2S0 등): 결합 occupancy는 공정 KPI가 아님. KPI는 catalytic competence",
        "• 화학공학 비유: 원료 스펙 시트(서열)만 보지 말고, 단위공정 흐름도 + 열역학 + 품질규격을 동시에 읽는다",
        "• 이 프로젝트의 1차 산출물 = 치료제가 아니라  ‘어느 결함 모드에 어떤 개입 논리를 시도할 가치가 있는가’ 의 공정 triage 맵",
    ], size=13, color=DARK, spacing=6)
    fin(s, "01 Why")

    # ══════════════════════════════════════════════════════════════════════
    # 4. CLINICAL / SCIENTIFIC CONTEXT
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "01 · 임상·과학 맥락 (사실 앵커)", "과장 없이 고정할 수 있는 외부 사실")
    facts = [
        ("질병", "MPS II (Hunter)\nIDS 결핍 → GAG 축적"),
        ("표적 효소", "Iduronate-2-sulfatase\nP22304 · sulfatase family"),
        ("실측 구조", "PDB 5FQL · 2.3 Å\nDemydchuk 2017"),
        ("변이 원천", "공개 ClinVar/UniProt\n프로젝트 병합 797건"),
        ("PC 선례", "migalastat (α-Gal A)\nambroxol–GCase 연구"),
        ("IDS PC 단서", "Hoshina 2018 등\n헤파린 유도체 계열"),
    ]
    for i, (t, d) in enumerate(facts):
        col, row = i % 3, i // 3
        left = Inches(0.4 + col * 4.25)
        top = Inches(1.3 + row * 2.7)
        card(s, left, top, Inches(4.0), Inches(2.4), fill=CARD, line=LINE)
        add_rect(s, left, top, Inches(4.0), Inches(0.5), fill=NAVY)
        add_text(s, left + Inches(0.2), top + Inches(0.1), Inches(3.6), Inches(0.35),
                 t, size=14, bold=True, color=MINT)
        add_text(s, left + Inches(0.25), top + Inches(0.75), Inches(3.5), Inches(1.4),
                 d, size=15, color=DARK)
    fin(s, "01 Context")

    # ══════════════════════════════════════════════════════════════════════
    # 5. DEVELOPMENT TIMELINE
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "02 · 개발 연대기 — 실제로 무엇이 쌓였는가", "시간순 · 산출물 중심 (요약을 과도하게 압축하지 않음)")
    timeline = [
        ("Phase 0", "고정", "5FQL.pdb 확보\nP22304 표적 고정\n‘도구 모음’ 유혹 거부"),
        ("Phase 1", "DB", "797 변이 병합\nRSA·촉매거리·Ca\n이황화·글리코 주석"),
        ("Phase 2", "Shortlist", "코어 불안정 필터\n안전권≥6.5Å\nSer152 등 hotspot"),
        ("Phase 3", "규칙 v1", "빈도→amenability\nHard EXCLUDE\nTIER1 proxy 점수"),
        ("Phase 4", "기전 심화", "Ser152 모듈 전수\nH-bond·패킹 순위\n환자용 파이프라인"),
        ("Phase 5", "AF Server", "WT vs S152I\nRMSD≈0.37Å\npLDDT 모듈↓"),
        ("Phase 6", "열역학", "EvoEF2 ΔΔG\nC422F·C184W\nCTSS·neo-pocket"),
        ("Phase 7", "공개층", "idsx-lab 포털\n자문 패키지\nclaim/non-claim"),
    ]
    for i, (ph, tag, body) in enumerate(timeline):
        col = i % 4
        row = i // 4
        left = Inches(0.35 + col * 3.25)
        top = Inches(1.25 + row * 2.85)
        card(s, left, top, Inches(3.1), Inches(2.65), fill=CARD, line=LINE)
        add_rect(s, left, top, Inches(3.1), Inches(0.55), fill=TEAL if row == 0 else NAVY)
        add_text(s, left + Inches(0.15), top + Inches(0.12), Inches(1.8), Inches(0.35),
                 ph, size=13, bold=True, color=WHITE)
        add_text(s, left + Inches(1.9), top + Inches(0.12), Inches(1.0), Inches(0.35),
                 tag, size=12, bold=True, color=MINT, align=PP_ALIGN.RIGHT)
        add_text(s, left + Inches(0.2), top + Inches(0.75), Inches(2.7), Inches(1.7),
                 body, size=13, color=DARK)
    fin(s, "02 Timeline")

    # ══════════════════════════════════════════════════════════════════════
    # 6. TRAPS AVOIDED
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "02 · 초기에 피하려고 한 함정 (사고의 선택)", "버린 길이 이 프로젝트의 정체성")
    traps = [
        ("유혹", "왜 위험한가", "택한 대안"),
        ("AI 설계 플랫폼부터", "질문 없는 도구 수집", "단일 단백질·단일 질병 고정"),
        ("RFdiffusion 먼저", "예쁜 복합체 ≠ 생물학", "변이–구조 분류·규칙 선행"),
        ("암브록솔 IDS 직적용", "표적 효소가 다름", "성공 조건만 이식(레퍼런스)"),
        ("병원성 빈도=우선순위", "null·촉매파괴도 빈도↑", "Hard EXCLUDE + 점수"),
        ("Ser152에 바인더 부착", "RSA 0% = 결합면 아님", "인접 표면 에피토프 원칙"),
        ("AF RMSD로 순위", "MSA 합의 fold 착시", "ΔΔG·구조 주석·실험 축"),
        ("열안정=프로젝트 제목", "온도계≠최종 KPI", "competence · ER 통과"),
    ]
    # header
    y0 = Inches(1.2)
    widths = [Inches(3.5), Inches(4.3), Inches(4.5)]
    xs = [Inches(0.4), Inches(3.9), Inches(8.2)]
    for j, htxt in enumerate(traps[0]):
        add_rect(s, xs[j], y0, widths[j], Inches(0.4), fill=NAVY)
        add_text(s, xs[j] + Inches(0.1), y0 + Inches(0.05), widths[j] - Inches(0.15), Inches(0.3),
                 htxt, size=12, bold=True, color=MINT)
    for i, row in enumerate(traps[1:]):
        y = y0 + Inches(0.45) + i * Inches(0.72)
        bg = OFF_WHITE if i % 2 == 0 else CARD
        for j, cell in enumerate(row):
            add_rect(s, xs[j], y, widths[j], Inches(0.68), fill=bg, line=LINE)
            add_text(s, xs[j] + Inches(0.1), y + Inches(0.15), widths[j] - Inches(0.15), Inches(0.45),
                     cell, size=12, bold=(j == 0), color=DARK if j else TEAL_DK)
    fin(s, "02 Traps")

    # ══════════════════════════════════════════════════════════════════════
    # 7. LOGIC 4-STAGE
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "03 · 로직의 뼈대 — 4단 추론", "이 네 문장이 프로젝트 전체를 지탱한다")
    steps = [
        ("①", "네트워크", "효소는 잔기 집합이 아니라\n구조 상호작용 네트워크 위에서 기능한다", TEAL),
        ("②", "결함 모드", "일부 미스센스는 활성부위를 부수기보다\n코어 패킹·H-bond 네트워크를 흔든다", PURPLE),
        ("③", "세포 경로", "접힘·ERQC·배달이 깨지면\n측정 활성은 분자 수 손실만으로도 급감", AMBER),
        ("④", "개입 논리", "고장 노드에 직접 붙이지 말고\n연결된 표면을 안정화해 competence를 본다", CORAL),
    ]
    for i, (n, t, d, c) in enumerate(steps):
        top = Inches(1.25 + i * 1.35)
        add_shape(s, Inches(0.4), top, Inches(1.1), Inches(1.15), fill=c)
        add_text(s, Inches(0.4), top + Inches(0.35), Inches(1.1), Inches(0.45),
                 n, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        card(s, Inches(1.7), top, Inches(11.2), Inches(1.15), fill=CARD, line=LINE)
        add_text(s, Inches(1.95), top + Inches(0.15), Inches(2.2), Inches(0.35),
                 t, size=16, bold=True, color=c)
        add_text(s, Inches(4.3), top + Inches(0.2), Inches(8.3), Inches(0.8),
                 d, size=14, color=DARK)
    fin(s, "03 Logic")

    # ══════════════════════════════════════════════════════════════════════
    # 8. TERMINOLOGY
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "03 · 용어 규율 — 여기서 헷갈리면 전체가 무너짐", "Hotspot · Module · Epitope · Amenability · Competence")
    terms = [
        ("Hotspot / TIER1", "구조적으로 아프고 구제 논리를 시도할 가치가 있는 변이 노드", "결합 자리 · ‘약 된다’ 증명"),
        ("Structural module", "그 노드 + H-bond·패킹으로 묶인 이웃 집합", "단일 아미노산 숭배"),
        ("Epitope / 표면 패치", "바인더가 닿을 후보 표면(모듈의 손잡이)", "변이 잔기 자체 · 임상 항원 필수 의미"),
        ("Amenability proxy", "구조 규칙 점수 (사전 등록 필터)", "잔존 활성 · 임상 반응"),
        ("Catalytic competence", "기능 구조 회복 후 실제 촉매 능력", "결합 occupancy · 예쁜 docking"),
        ("열적 안정성 (Tm)", "네트워크 유지의 대리 지표(온도계)", "프로젝트 제목 · 최종 목표"),
    ]
    for i, (name, is_, isnot) in enumerate(terms):
        y = Inches(1.2 + i * 0.9)
        add_rect(s, Inches(0.35), y, Inches(2.8), Inches(0.8), fill=NAVY)
        add_text(s, Inches(0.45), y + Inches(0.22), Inches(2.6), Inches(0.45),
                 name, size=12, bold=True, color=WHITE)
        card(s, Inches(3.25), y, Inches(5.0), Inches(0.8), fill=SOFT_GREEN, line=GREEN)
        add_text(s, Inches(3.4), y + Inches(0.08), Inches(4.7), Inches(0.25),
                 "의미 (IDSX)", size=10, bold=True, color=GREEN)
        add_text(s, Inches(3.4), y + Inches(0.35), Inches(4.7), Inches(0.4),
                 is_, size=12, color=DARK)
        card(s, Inches(8.4), y, Inches(4.5), Inches(0.8), fill=SOFT_CORAL, line=CORAL)
        add_text(s, Inches(8.55), y + Inches(0.08), Inches(4.2), Inches(0.25),
                 "아닌 것", size=10, bold=True, color=CORAL)
        add_text(s, Inches(8.55), y + Inches(0.35), Inches(4.2), Inches(0.4),
                 isnot, size=12, color=DARK)
    fin(s, "03 Terms")

    # ══════════════════════════════════════════════════════════════════════
    # 9. SUCCESS HIERARCHY
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "03 · 목표 계층도 — 무엇이 목표이고 무엇이 온도계인가",
                "암브록솔 논의에서 확정한 구분")
    levels = [
        (Inches(0.5), "최종 목표 (Biology)", "리소좀에서 작동하는 IDS 분자 수 · 활성\n= catalytic competence", CORAL, SOFT_CORAL),
        (Inches(2.0), "세포 중간 목표", "ERQC 통과 · 골지 성숙 · 리소좀 배달 ↑  /  ERAD ↓", AMBER, SOFT_AMBER),
        (Inches(3.5), "분자 중간 목표", "구조 모듈·interaction network가 WT-like 접힘을 더 잘 유지", TEAL, SOFT_TEAL),
        (Inches(5.0), "실험실 대리 지표 예", "Tm · ΔΔG · 프로테아제 민감도 · 성숙 밴드 · KD · 국재", PURPLE, SOFT_PURPLE),
    ]
    for top, title, body, ac, bg in levels:
        card(s, Inches(1.5), top, Inches(10.3), Inches(1.25), fill=bg, line=ac)
        add_text(s, Inches(1.75), top + Inches(0.15), Inches(9.8), Inches(0.35),
                 title, size=15, bold=True, color=ac)
        add_text(s, Inches(1.75), top + Inches(0.55), Inches(9.8), Inches(0.55),
                 body, size=14, color=DARK)
        if top != Inches(5.0):
            add_text(s, Inches(6.3), top + Inches(1.15), Inches(0.5), Inches(0.3),
                     "↑", size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    fin(s, "03 Hierarchy")

    # ══════════════════════════════════════════════════════════════════════
    # 10. WHAT WAS BUILT — ARTIFACTS
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "04 · 현재까지 구축된 산출물 지도", "파일은 흩어져 있으나 논리는 한 줄기")
    arts = [
        ("변이–구조 DB", "IDS_P22304_Mutation_\nStructural_Database.xlsx", "797행 · 6시트\nRSA/거리/결함범주"),
        ("Binder shortlist", "IDS_Chaperone_Binder_\nTarget_Shortlist.xlsx", "안전권 15 · 주의 15\n에피토프 후보"),
        ("Amenability v1", "Rules + Prototype xlsx\nClaim/Ambroxol 문서", "EXCLUDE/CAUTION/TIER\n전량 점수 적용"),
        ("Ser152 심화", "Mechanism Review\nFig1–4 + 체크리스트", "기전 순위 · 2축 맵\n환자용 파이프라인"),
        ("AF 계산", "fold_WT / S152I / multi\nAF3 분석 보고서", "RMSD 0.37Å\n모듈 pLDDT↓"),
        ("열역학 층", "EvoEF2 ΔΔG · CTSS\nc422f pocket", "에너지 vs 노출\nneo-pocket 가설"),
        ("자문 패키지", "Advisor_Briefing_Package\n질문지 · 브리핑", "한페이지 · 궤적\n면담 질문"),
        ("공개 포털", "idsx-lab (GitHub Pages)\nframework·compute·AF", "claim 고정\n사고 타임라인"),
    ]
    for i, (t, f, d) in enumerate(arts):
        col, row = i % 4, i // 4
        left = Inches(0.35 + col * 3.25)
        top = Inches(1.25 + row * 2.85)
        card(s, left, top, Inches(3.1), Inches(2.65), fill=CARD, line=LINE)
        add_rect(s, left, top, Inches(3.1), Inches(0.5), fill=TEAL)
        add_text(s, left + Inches(0.15), top + Inches(0.1), Inches(2.8), Inches(0.35),
                 t, size=13, bold=True, color=WHITE)
        add_text(s, left + Inches(0.15), top + Inches(0.7), Inches(2.8), Inches(0.9),
                 f, size=11, color=TEAL_DK)
        add_text(s, left + Inches(0.15), top + Inches(1.7), Inches(2.8), Inches(0.75),
                 d, size=12, color=DARK)
    fin(s, "04 Artifacts")

    # ══════════════════════════════════════════════════════════════════════
    # 11. AMENABILITY PIPELINE
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "04 · Amenability 파이프라인 (규칙의 흐름)", "빈도 랭킹을 버리고 개별 변이 점수로 재편")
    funnel = [
        ("입력", "ClinVar/UniProt × 5FQL\n797 variants", NAVY),
        ("Hard EXCLUDE", "null · 촉매 · Ca\n이황화 · 미해석", CORAL),
        ("CAUTION", "촉매거리 < 6.5 Å\nArg88 계열", AMBER),
        ("점수 0–10", "코어/라벨/이격\nRSA/에피토프", TEAL),
        ("TIER1 proxy", "score ≥ 7\n위치 collapse 필요", GREEN),
        ("모달리티", "A 소분자 PC\nB 표면 바인더\nC 비구제", PURPLE),
    ]
    for i, (t, d, c) in enumerate(funnel):
        left = Inches(0.3 + i * 2.15)
        card(s, left, Inches(1.35), Inches(2.05), Inches(2.6), fill=CARD, line=c)
        add_rect(s, left, Inches(1.35), Inches(2.05), Inches(0.45), fill=c)
        add_text(s, left + Inches(0.05), Inches(1.42), Inches(1.95), Inches(0.35),
                 t, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, left + Inches(0.1), Inches(2.0), Inches(1.85), Inches(1.7),
                 d, size=12, color=DARK, align=PP_ALIGN.CENTER)
        if i < 5:
            add_text(s, left + Inches(1.9), Inches(2.4), Inches(0.3), Inches(0.4),
                     "›", size=20, bold=True, color=NAVY)

    # counts
    card(s, Inches(0.4), Inches(4.2), Inches(12.5), Inches(2.6), fill=LIGHT_BG, line=LINE)
    add_text(s, Inches(0.65), Inches(4.35), Inches(12), Inches(0.35),
             "v1.0 규칙 전량 러프 적용 결과 (경향 확인용 · 1점 단위 차 가능)", size=13, bold=True, color=NAVY)
    counts = [("EXCLUDE", "~210–231"), ("CAUTION_POCKET", "~106"), ("TIER1 proxy", "~146–163"),
              ("TIER2", "~227–244"), ("NEG_CTRL", "~64–85"), ("TIER3", "~6")]
    for i, (k, v) in enumerate(counts):
        left = Inches(0.65 + i * 2.05)
        add_text(s, left, Inches(4.9), Inches(1.9), Inches(0.35), k, size=11, bold=True, color=TEAL)
        add_text(s, left, Inches(5.3), Inches(1.9), Inches(0.4), v, size=18, bold=True, color=NAVY)
    add_text(s, Inches(0.65), Inches(5.9), Inches(12), Inches(0.6),
             "해석 규율: TIER1 n은 ‘치료 후보 n개’가 아니다. 위치 단위 collapse + 잔존활성 문헌으로 2차 축소. "
             "1편 본문 강조는 위치 15–30 + 심화 2–3개(Ser152, Asp308, Tyr108)가 적절.",
             size=12, color=GRAY)
    fin(s, "04 Pipeline")

    # ══════════════════════════════════════════════════════════════════════
    # 12. SER152 CASE
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "05 · 대표 케이스 Ser152 — 구조 모듈 (얼굴 케이스)",
                "촉매 파괴가 아니라 코어 핀 · 개입은 인접 표면")
    # left numbers
    card(s, Inches(0.35), Inches(1.2), Inches(6.3), Inches(5.7), fill=CARD, line=LINE)
    add_text(s, Inches(0.55), Inches(1.35), Inches(5.9), Inches(0.35),
             "5FQL 실측 앵커 [S]", size=14, bold=True, color=TEAL)
    rows = [
        ("RSA", "0% (완전 매장)"),
        ("촉매 포켓 거리", "10.27 Å (안전권 ≥6.5)"),
        ("H-bond (Ser-OH)", "Asp148 ~2.53 Å · Thr118 ~2.85 Å"),
        ("부피 변화 Ser→Ile", "≈ +77.7 Å³"),
        ("소수성 변화 (KD)", "+5.3 (극성 틈 → 지방족)"),
        ("패킹 (CB 4.5Å)", "heavy atom 21개 — 조밀"),
        ("임상 라벨", "Pathogenic 계열 반복 (S152I/R/N 등)"),
        ("IDSX 분류", "Core destabilization · score 10 · TIER1"),
        ("결합 전략", "152 금지 · Tyr151/Pro150/Ser154 패치"),
        ("1차 기전 순위", "부피/패킹 > H-bond 상실 > 극성 불일치"),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(1.8 + i * 0.48)
        add_text(s, Inches(0.55), y, Inches(2.3), Inches(0.4), k, size=12, bold=True, color=GRAY)
        add_text(s, Inches(2.9), y, Inches(3.5), Inches(0.4), v, size=12, color=DARK)

    # right patient analogy
    card(s, Inches(6.85), Inches(1.2), Inches(6.1), Inches(2.6), fill=SOFT_TEAL, line=TEAL)
    add_text(s, Inches(7.1), Inches(1.35), Inches(5.6), Inches(0.35),
             "공정 비유 (화학공학)", size=14, bold=True, color=TEAL_DK)
    add_text(s, Inches(7.1), Inches(1.8), Inches(5.6), Inches(1.8),
             "리소좀 = 현장 반응기, ER = 조립 공장.\n"
             "S152I는 칼날(활성부위)을 직접 부러뜨리기보다\n"
             "조립 핀(코어 패킹·H-bond)이 어긋나 QC에서 폐기·\n"
             "미배달되는 유형에 가깝다고 읽을 구조적 이유가 있다.\n"
             "(개인 잔존활성 %는 미측정 — UNKNOWN)",
             size=13, color=DARK)

    card(s, Inches(6.85), Inches(4.0), Inches(6.1), Inches(2.9), fill=SOFT_AMBER, line=AMBER)
    add_text(s, Inches(7.1), Inches(4.15), Inches(5.6), Inches(0.35),
             "기전 체인 (측정 활성 ↓)", size=14, bold=True, color=AMBER)
    add_text(s, Inches(7.1), Inches(4.6), Inches(5.6), Inches(2.1),
             "측쇄 물리 결함 (부피+H-bond+극성)\n"
             "→ 접힘 안정성↓ / 국소 풀림 경향↑\n"
             "→ ERQC 지연·불합격 → ERAD·배달 실패\n"
             "→ 리소좀 기능 분자 수↓ → assay 활성↓\n"
             "※ 클래스 문헌 지지 · 이 대립유전자 단독 실측 약함",
             size=13, color=DARK)
    fin(s, "05 Ser152")

    # ══════════════════════════════════════════════════════════════════════
    # 13. SER152 FIGURE
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "05 · Ser152 기전 시각 자료", "Mechanism Review Fig — 기록·자문용")
    if not try_pic(s, IMG["cascade"], Inches(0.3), Inches(1.15), Inches(6.3), Inches(5.8)):
        card(s, Inches(0.3), Inches(1.15), Inches(6.3), Inches(5.8), fill=LIGHT_BG)
        add_text(s, Inches(0.5), Inches(3.5), Inches(5.9), Inches(0.5),
                 "Fig1 cascade (파일 경로 확인)", size=14, color=GRAY, align=PP_ALIGN.CENTER)
    if not try_pic(s, IMG["2axis"], Inches(6.7), Inches(1.15), Inches(6.2), Inches(5.8)):
        card(s, Inches(6.7), Inches(1.15), Inches(6.2), Inches(5.8), fill=LIGHT_BG)
    fin(s, "05 Figs")

    # ══════════════════════════════════════════════════════════════════════
    # 14. ARG88 COUNTEREXAMPLE
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "05 · 교육용 반례 Arg88 — 빈도 ≠ 적합성", "이 figure가 없으면 triage가 ‘병원성 랭킹’으로 퇴화한다")
    card(s, Inches(0.4), Inches(1.25), Inches(6.1), Inches(5.6), fill=SOFT_CORAL, line=CORAL)
    add_text(s, Inches(0.65), Inches(1.45), Inches(5.6), Inches(0.4),
             "Arg88 프로필", size=16, bold=True, color=CORAL)
    add_paras(s, Inches(0.65), Inches(2.0), Inches(5.6), Inches(4.5), [
        "• 병원성 미스센스 반복 빈도: 높음 (shortlist 상위권)",
        "• 촉매 포켓 거리: 2.77 Å  —  극도 인접",
        "• 분류: CAUTION_NEAR_POCKET (바인더 1순위 금지)",
        "• 위험: 결합제 접근 시 기질 채널·활성 미세환경 간섭",
        "• 교훈: ‘많이 아프다고 보고된 자리’ ≠ ‘밖에서 붙잡아 구제할 자리’",
        "• 대안 모달리티: 활성부위형 소분자 inhibitory PC는 별도 트랙",
        "  (억제–리소좀 해리 창이 필수 — 학부 범위에서는 개념 층)",
        "• 논문 figure 역할: 빈도 중심 shortlist의 실패 모드를 스스로 고발",
    ], size=13, color=DARK, spacing=7)

    card(s, Inches(6.8), Inches(1.25), Inches(6.1), Inches(5.6), fill=SOFT_GREEN, line=GREEN)
    add_text(s, Inches(7.05), Inches(1.45), Inches(5.6), Inches(0.4),
             "대비 표 — 우선 착수 후보", size=16, bold=True, color=GREEN)
    add_paras(s, Inches(7.05), Inches(2.0), Inches(5.6), Inches(4.5), [
        "Ser152   | 병원성 다수 | dist 10.27 | RSA 0% | 에피토프 11 | 1순위",
        "Asp308   | 병원성 다수 | dist 10.87 | 매장 | 에피토프 6  | 1순위",
        "Tyr108   | 혼합       | dist 8.11  | 경계 | 에피토프 14 | 1순위",
        "Pro469   | 중         | dist 12.02 | —    | 에피토프 10 | 2순위",
        "Arg468   | 높음       | dist 6.03  | 경계 | —           | 재검토",
        "Arg88    | 높음       | dist 2.77  | —    | —           | 주의/배제",
        "",
        "원칙: 안전권(≥6.5Å)에서 구조 모듈 + 표면 손잡이가 있는 자리부터.",
        "주의: 6.5Å은 마법의 숫자가 아님 → 자문 질문 1순위.",
    ], size=13, color=DARK, spacing=6)
    fin(s, "05 Arg88")

    # ══════════════════════════════════════════════════════════════════════
    # 15. AF RESULTS
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "04 · AlphaFold Server — WT vs S152I (해석 규율 포함)",
                "좌표가 비슷하다고 해서 안정·무해가 아니다")
    card(s, Inches(0.35), Inches(1.2), Inches(6.4), Inches(3.3), fill=CARD, line=LINE)
    add_text(s, Inches(0.55), Inches(1.35), Inches(6), Inches(0.35),
             "핵심 수치 (model_0, UniProt 34–550)", size=13, bold=True, color=TEAL)
    metrics = [
        ("Global Cα RMSD", "0.372 Å"),
        ("Local U145–155 RMSD", "0.274 Å"),
        ("pTM WT / S152I", "0.97 / 0.96"),
        ("pLDDT U152", "98.8 → 84.2"),
        ("pLDDT U110–160 mean", "98.15 → 91.38"),
        ("Cα dev at U152", "0.909 Å"),
    ]
    for i, (k, v) in enumerate(metrics):
        y = Inches(1.85 + i * 0.4)
        add_text(s, Inches(0.55), y, Inches(3.5), Inches(0.35), k, size=13, color=GRAY)
        add_text(s, Inches(4.0), y, Inches(2.5), Inches(0.35), v, size=13, bold=True, color=NAVY)

    card(s, Inches(6.95), Inches(1.2), Inches(5.95), Inches(3.3), fill=SOFT_AMBER, line=AMBER)
    add_text(s, Inches(7.15), Inches(1.35), Inches(5.5), Inches(0.35),
             "해석 — 지지 / 기각불가 / 금지", size=13, bold=True, color=AMBER)
    add_text(s, Inches(7.15), Inches(1.85), Inches(5.5), Inches(2.4),
             "지지: 접힌 native-like 좌표는 양쪽에 존재 가능.\n"
             "대붕괴 misfolding 전제는 약화.\n\n"
             "기각 못함: ERAD · ΔΔG · 잔존활성 · 배달.\n\n"
             "금지: RMSD 작다=무해 · pLDDT 높다=치료불필요.\n\n"
             "다음: template-off 교차 · FoldX/Rosetta ΔΔG.",
             size=12, color=DARK)

    # AF paradigm
    card(s, Inches(0.35), Inches(4.7), Inches(12.55), Inches(2.15), fill=SOFT_PURPLE, line=PURPLE)
    add_text(s, Inches(0.55), Inches(4.85), Inches(12.2), Inches(0.35),
             "패러다임 전환 (본 연구자의 핵심 통찰)", size=13, bold=True, color=PURPLE)
    add_text(s, Inches(0.55), Inches(5.3), Inches(12.2), Inches(1.35),
             "AlphaFold는 물리학(ΔG)보다 MSA 진화 합의에 크게 의존한다. 점 돌연변이 하나가 MSA 전체를 바꾸지 않으므로\n"
             "내부 코어 충돌이 있어도 WT 뼈대에 억지 끼워 맞춘 ‘예쁜 그림’이 나올 수 있다. 따라서 hotspot 검증의 중심축을\n"
             "‘backbone RMSD 비교’에서 ‘구조 주석 + 열역학(ΔΔG) + 세포 전제’로 옮긴 것이 IDSX의 방법론적 정체성이다.\n"
             "이 전환은 실패가 아니라, 도구의 작동 원리를 이해한 뒤의 설계 선택이다.",
             size=13, color=DARK)
    fin(s, "04 AF")

    # ══════════════════════════════════════════════════════════════════════
    # 16. AF IMAGES
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "04 · AF 시각 자료", "오버레이 · 프로파일 · 모듈 줌")
    try_pic(s, IMG["af_overlay"], Inches(0.25), Inches(1.15), Inches(4.2), Inches(5.7))
    try_pic(s, IMG["af_profiles"], Inches(4.55), Inches(1.15), Inches(4.2), Inches(5.7))
    try_pic(s, IMG["af_zoom"], Inches(8.85), Inches(1.15), Inches(4.2), Inches(5.7))
    fin(s, "04 AF img")

    # ══════════════════════════════════════════════════════════════════════
    # 17. DDG + CTSS TENSION
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "04 · ΔΔG / CTSS 층 — 그리고 트랙 긴장 (정직하게)",
                "에너지 폭발 ≠ 약물 접근 가능. 두 트랙을 한 문장으로 섞지 않는다.")
    card(s, Inches(0.35), Inches(1.2), Inches(6.3), Inches(5.6), fill=CARD, line=LINE)
    add_text(s, Inches(0.55), Inches(1.35), Inches(5.9), Inches(0.35),
             "EvoEF2 ΔΔG (kcal/mol) 일부", size=14, bold=True, color=TEAL)
    ddg = [
        ("C422F", "+121.62", "이황화 파괴 + 거대 방향족"),
        ("C184W", "+111.32", "이황화 파괴 + Trp 삽입"),
        ("S152I", "+46.02", "코어 모듈 · 중간 고에너지"),
        ("L339P", "+38.97", "헬릭스 kink"),
        ("W345R", "+31.97", "코어 전하"),
        ("V503D", "+7.91", "코어 극성"),
    ]
    for i, (m, e, note) in enumerate(ddg):
        y = Inches(1.9 + i * 0.7)
        add_text(s, Inches(0.55), y, Inches(1.5), Inches(0.35), m, size=14, bold=True, color=NAVY)
        add_text(s, Inches(2.1), y, Inches(1.5), Inches(0.35), e, size=14, bold=True, color=CORAL)
        add_text(s, Inches(3.7), y, Inches(2.7), Inches(0.55), note, size=12, color=GRAY)

    card(s, Inches(6.85), Inches(1.2), Inches(6.1), Inches(5.6), fill=SOFT_AMBER, line=AMBER)
    add_text(s, Inches(7.1), Inches(1.35), Inches(5.6), Inches(0.4),
             "두 트랙의 긴장 — 숨기지 않는다", size=14, bold=True, color=AMBER)
    add_text(s, Inches(7.1), Inches(1.9), Inches(5.6), Inches(4.6),
             "Track A (Amenability / 바인더 논리)\n"
             "• Ser152 같은 매장 모듈: 구제 논리는 ‘인접 표면 패치’\n"
             "• CTSS 관점 SASA≈0 → ‘변이 잔기 직접 소분자’에는 undruggable\n"
             "• 그러나 알로스테릭 클램프는 여전히 가설로 남음\n\n"
             "Track B (Neo-pocket / 소분자 클램프)\n"
             "• C422F/Y: 고 ΔΔG + 표면 노출 소수성 틈 (CTSS 상위)\n"
             "• 파마코포어: 양극성 클램프 + 중심 방향족\n"
             "• Hard EXCLUDE 이황화 파괴형 → PC 반응 낮을 수 있음\n"
             "  (amenability 규칙과 충돌 가능 — 자문 필수)\n\n"
             "연구 원칙: 트랙을 섞어 ‘하나의 치료 스토리’로 포장하지 않는다.\n"
             "각각 다른 결함 모드 · 다른 모달리티 · 다른 성공 기준.",
             size=12, color=DARK)
    fin(s, "04 Tension")

    # ══════════════════════════════════════════════════════════════════════
    # 18. CTSS VISUAL
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "04 · CTSS 산점도 · C422F neo-pocket (Track B 자료)",
                "노출(SASA)×에너지×소수성 — 소분자 접근성 렌즈")
    try_pic(s, IMG["ctss"], Inches(0.3), Inches(1.15), Inches(6.5), Inches(5.7))
    try_pic(s, IMG["c422"], Inches(6.95), Inches(1.15), Inches(6.0), Inches(5.7))
    fin(s, "04 CTSS")

    # ══════════════════════════════════════════════════════════════════════
    # 19. STRENGTH MAP
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "06 · 로직 강도 맵 — 어디가 단단하고 어디가 열려 있는가",
                "자문·리뷰어에게 가장 정직한 슬라이드")
    strength = [
        ("네트워크·competence 프레임", "★★★★☆", "방법론적으로 탄탄", TEAL),
        ("Ser152 구조 모듈 서술", "★★★★☆", "5FQL H-bond·RSA·거리 뒷받침", TEAL),
        ("Amenability triage 규칙", "★★★☆☆", "재현 가능 · 컷오프는 자문 대상", AMBER),
        ("에피토프 간접 안정화 전략", "★★★☆☆", "논리 일관 · 미검증", AMBER),
        ("AF 해석 규율", "★★★★☆", "과대/과소 방지 장치 명시", TEAL),
        ("ΔΔG/CTSS 정량", "★★★☆☆", "유용 렌즈 · forcefield 한계·EXCLUDE 긴장", AMBER),
        ("세포·잔존활성 전제", "★★☆☆☆", "클래스 문헌 · 대립유전자 실측 약함", CORAL),
        ("치료·구제 성공", "★☆☆☆☆", "claim 금지 영역", CORAL),
    ]
    for i, (name, stars, note, c) in enumerate(strength):
        y = Inches(1.2 + i * 0.7)
        card(s, Inches(0.4), y, Inches(12.5), Inches(0.62), fill=CARD, line=LINE)
        add_text(s, Inches(0.6), y + Inches(0.15), Inches(5.5), Inches(0.35), name, size=13, bold=True, color=DARK)
        add_text(s, Inches(6.3), y + Inches(0.15), Inches(2.0), Inches(0.35), stars, size=14, bold=True, color=c)
        add_text(s, Inches(8.4), y + Inches(0.15), Inches(4.3), Inches(0.35), note, size=12, color=GRAY)
    fin(s, "06 Strength")

    # ══════════════════════════════════════════════════════════════════════
    # 20. SIGNIFICANCE
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "06 · 이 프로젝트의 유의미성 (왜 계속할 가치가 있는가)",
                "치료 과대선전이 아니라 과학·교육·진로 층에서")
    sigs = [
        ("과학", "LSD PC 전략에서 가장 흔한 실패 모드 — ‘누구에게 시도할 것인가’ — 를 구조 규칙으로 사전 분리한다. "
         "도구(RFdiffusion 등) 전에 질문과 필터를 고정하는 연구 문화 자체에 기여."),
        ("방법", "MSA-AI 착시 → 열역학·구조 주석으로 패러다임 전환을 문서화. 학부 수준에서도 재현 가능한 "
         "triage 프로토콜(규칙 사전등록·non-claim·에피토프 분리)을 제시."),
        ("임상 가교", "환자·가족이 오해하기 쉬운 ‘칼날이 없다 vs 배달이 안 된다’를 구조 언어로 구분. "
         "개인 치료 권고는 금지하되, 대화의 정확성을 올린다."),
        ("화학공학", "공정 흐름도·결함 모드·KPI 계층·에너지 수지 사고를 단백질 질병에 이식. "
         "전공 정체성이 반영된 독자적 프레임."),
        ("진로·협업", "구조생물·단백질공학·리소좀 임상의학 랩에 가져갈 수 있는 ‘완성된 질문 패키지’. "
         "결과 자랑보다 critique를 요청하는 형식."),
        ("공개 기록", "idsx-lab 포털로 claim/non-claim과 사고 궤적을 외부에 고정. "
         "혼자 규칙을 유리하게 바꾸지 못하도록 자기 구속."),
    ]
    for i, (t, d) in enumerate(sigs):
        col, row = i % 3, i // 3
        left = Inches(0.35 + col * 4.3)
        top = Inches(1.25 + row * 2.85)
        card(s, left, top, Inches(4.1), Inches(2.65), fill=CARD, line=LINE)
        add_rect(s, left, top, Inches(4.1), Inches(0.5), fill=NAVY)
        add_text(s, left + Inches(0.2), top + Inches(0.1), Inches(3.7), Inches(0.35),
                 t, size=14, bold=True, color=MINT)
        add_text(s, left + Inches(0.2), top + Inches(0.7), Inches(3.7), Inches(1.8),
                 d, size=12, color=DARK)
    fin(s, "06 Significance")

    # ══════════════════════════════════════════════════════════════════════
    # 21. RISKS
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "06 · 리스크 레지스터 (알고 있는 실패 가능 지점)", "숨기지 않는 것이 전문가 신뢰")
    risks = [
        ("R1", "ClinVar 라벨 노이즈", "병원성 가점이 틀림 → 잔존활성 문헌으로 교정"),
        ("R2", "Core destabil 규칙 ≠ 실측 ΔΔG", "FoldX/Rosetta/실험으로 교차 검증"),
        ("R3", "에피토프 개수 ≠ 패치 품질", "표면 형상·소수성·요철 사전 스크리닝"),
        ("R4", "리소좀 pH·프로테아제", "바인더 스캐폴드 안정성 엔지니어링 필요"),
        ("R5", "CNS 전달 (신경형)", "치료 overclaim 금지 · 별도 전달 문제"),
        ("R6", "계산만 쌓임", "랩 소속·최소 assay 없이 리뷰어 불신"),
        ("R7", "Track A/B 혼동", "모달리티·EXCLUDE 규칙 충돌을 명시적 관리"),
        ("R8", "6.5Å 자의성", "자문·민감도 분석으로 작업 컷오프 지위 유지"),
    ]
    for i, (code, title, action) in enumerate(risks):
        col, row = i % 2, i // 2
        left = Inches(0.35 + col * 6.5)
        top = Inches(1.2 + row * 1.4)
        card(s, left, top, Inches(6.3), Inches(1.25), fill=CARD, line=LINE)
        add_shape(s, left + Inches(0.15), top + Inches(0.3), Inches(0.9), Inches(0.65), fill=CORAL)
        add_text(s, left + Inches(0.15), top + Inches(0.42), Inches(0.9), Inches(0.4),
                 code, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, left + Inches(1.25), top + Inches(0.2), Inches(4.8), Inches(0.35),
                 title, size=14, bold=True, color=NAVY)
        add_text(s, left + Inches(1.25), top + Inches(0.6), Inches(4.8), Inches(0.5),
                 action, size=12, color=GRAY)
    fin(s, "06 Risks")

    # ══════════════════════════════════════════════════════════════════════
    # 22. 90-DAY PLAN
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "07 · 향후 90일 계획 (우선순위 고정안)", "한 갈래만 깊게 — 도구 수집 금지")
    phases = [
        ("D0–14", "재료 정합", [
            "Ser152 계열 ClinVar_id/HGVS 정리",
            "말투 통일 (최중요 잔기 금지)",
            "Track A/B 문서 경계 문단 삽입",
            "자문 패키지 + 본 덱 송부",
        ], TEAL),
        ("D15–45", "숫자 층", [
            "FoldX 또는 Rosetta ΔΔG: WT vs S152I/R/N",
            "대조: Arg88 또는 표면 benign",
            "가능 시 AF template-off 1회",
            "CTSS 상위 C422의 amenability 재검토",
        ], PURPLE),
        ("D46–70", "생물 전제", [
            "TIER1 잔존활성·표현형 문헌 표",
            "UNKNOWN 채우기 / 채울 수 없음 명시",
            "IDS PC 문헌(Hoshina 등) 재독",
            "구제=가설 / triage=우선 유지",
        ], AMBER),
        ("D71–90", "설계 전 게이트", [
            "표면 패치 2–3개 시각 확정",
            "성공 기준 사전 등록 (발현/성숙/활성)",
            "한 트랙만 파일럿 여부 결정",
            "분류 논문 목차 초안",
        ], CORAL),
    ]
    for i, (when, title, bullets, c) in enumerate(phases):
        left = Inches(0.3 + i * 3.25)
        card(s, left, Inches(1.25), Inches(3.1), Inches(5.6), fill=CARD, line=c)
        add_rect(s, left, Inches(1.25), Inches(3.1), Inches(0.9), fill=c)
        add_text(s, left + Inches(0.15), Inches(1.35), Inches(2.8), Inches(0.3),
                 when, size=12, bold=True, color=WHITE)
        add_text(s, left + Inches(0.15), Inches(1.7), Inches(2.8), Inches(0.35),
                 title, size=16, bold=True, color=WHITE)
        add_paras(s, left + Inches(0.2), Inches(2.4), Inches(2.7), Inches(4.2),
                  ["• " + b for b in bullets], size=12, color=DARK, spacing=10)
    fin(s, "07 90d")

    # ══════════════════════════════════════════════════════════════════════
    # 23. 12-MONTH PLAN
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "07 · 12개월 로드맵 (게이트 포함)", "각 분기 끝에서 go/no-go")
    quarters = [
        ("Q1", "지도 완성", "triage 맵 고정\nΔΔG 패널\n잔존활성 문헌표\n자문 1–2회",
         "Gate: 본문 강조 위치 ≤20개로 축소"),
        ("Q2", "파일럿 1트랙", "패치 1개 확정\n소분자 OR 바인더\n계산 파이프만 또는\n협업 assay 설계",
         "Gate: 성공 기준 사전등록 문서"),
        ("Q3", "검증 접촉", "발현/DSF/SPR 중\n가능한 최소 실험\n또는 공개 데이터 재분석\n초고 Methods 확정",
         "Gate: non-claim 위반 0건 리뷰"),
        ("Q4", "산출·공개", "분류 논문/보고서\n데이터·코드 정리\nidsx-lab 업데이트\n랩 합류 포트폴리오",
         "Gate: 외부 1명 critique 반영"),
    ]
    for i, (q, title, body, gate) in enumerate(quarters):
        left = Inches(0.35 + i * 3.25)
        card(s, left, Inches(1.25), Inches(3.1), Inches(5.6), fill=CARD, line=LINE)
        add_rect(s, left, Inches(1.25), Inches(3.1), Inches(0.7), fill=NAVY)
        add_text(s, left + Inches(0.15), Inches(1.35), Inches(1.0), Inches(0.5),
                 q, size=20, bold=True, color=MINT)
        add_text(s, left + Inches(1.2), Inches(1.45), Inches(1.7), Inches(0.4),
                 title, size=14, bold=True, color=WHITE)
        add_text(s, left + Inches(0.2), Inches(2.2), Inches(2.7), Inches(2.5),
                 body, size=13, color=DARK)
        add_rect(s, left + Inches(0.15), Inches(5.0), Inches(2.8), Inches(1.5), fill=SOFT_TEAL)
        add_text(s, left + Inches(0.25), Inches(5.15), Inches(2.6), Inches(1.2),
                 gate, size=12, bold=True, color=TEAL_DK)
    fin(s, "07 12m")

    # ══════════════════════════════════════════════════════════════════════
    # 24. DO NOT DO
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "07 · 당장은 하지 말 것 (자기 규율)", "속도보다 방향")
    donts = [
        "전 변이 RFdiffusion 대량 생산 — 질문 없는 컴퓨팅",
        "암브록솔 IDS 직접 투여 claim",
        "AF RMSD만으로 논문 결론",
        "TIER1 목록을 환자 대상 치료 안내로 사용",
        "Track A와 Track B를 하나의 ‘최종 타겟 C422’로 강제 통합",
        "이황화 파괴형(C184/C422)을 amenability EXCLUDE와 충돌 없이 치료 후보로 포장",
        "열안정 보장 인자 찾기로 프로젝트 rename",
        "지도·윤리 없이 bioRxiv 성급한 투고",
        "신경 전달 미해결 상태에서 치료 완결 서사",
        "새 도구를 쌓기 위해 Step A–D 순서를 건너뛰기",
    ]
    for i, d in enumerate(donts):
        col, row = i % 2, i // 2
        left = Inches(0.35 + col * 6.5)
        top = Inches(1.2 + row * 1.05)
        card(s, left, top, Inches(6.3), Inches(0.95), fill=SOFT_CORAL, line=CORAL)
        add_text(s, left + Inches(0.2), top + Inches(0.15), Inches(0.5), Inches(0.6),
                 "✕", size=18, bold=True, color=CORAL)
        add_text(s, left + Inches(0.7), top + Inches(0.25), Inches(5.4), Inches(0.55),
                 d, size=13, color=DARK)
    fin(s, "07 Don't")

    # ══════════════════════════════════════════════════════════════════════
    # 25. ADVISOR QUESTIONS
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "07 · 자문에게 꼭 물을 것 (오늘 세 가지 + 확장)",
                "‘진행해도 되나’가 아니라 ‘어디를 고칠까’")
    card(s, Inches(0.4), Inches(1.25), Inches(12.5), Inches(2.0), fill=SOFT_TEAL, line=TEAL)
    add_text(s, Inches(0.65), Inches(1.4), Inches(12), Inches(0.35),
             "오늘 꼭 듣고 싶은 세 가지", size=14, bold=True, color=TEAL_DK)
    add_text(s, Inches(0.65), Inches(1.9), Inches(12), Inches(1.15),
             "1. 이 claim(5FQL 기반 amenability triage 맵)으로 계속 갈까, 방향을 틀까?\n"
             "2. Ser152 / Asp308 프로토타입이 첫 표·얼굴 케이스로 적절한가?\n"
             "3. 앞으로 90일 작업 순서 두 가지는? (E1 잔존활성 문헌 / E2 ΔΔG / E4 바인더 파일럿 / E6 랩 합류 중)",
             size=14, color=DARK)

    qs = [
        ("A", "1편 claim이 학문적으로 성립하는가? 실험 없이 구조 분류 논문의 가치는?"),
        ("B", "6.5Å 컷오프 · 이황화 EXCLUDE · ClinVar 가중 — 과도/부족/적절?"),
        ("C", "매장 hotspot + 인접 에피토프 프레임이 리뷰어에게 설득력 있는가?"),
        ("D", "GCase–ambroxol을 조건 표로만 쓰는 구성 · IDS PC 문헌 비중은?"),
        ("E", "Track B(C422 neo-pocket)를 병행할지, amenability 규칙과 어떻게 화해할지?"),
        ("F", "학부 단독 범위 · 붙을 랩 유형 · bioRxiv 시기 · 한 줄 숙제"),
    ]
    for i, (code, q) in enumerate(qs):
        y = Inches(3.5 + i * 0.55)
        add_shape(s, Inches(0.4), y, Inches(0.55), Inches(0.45), fill=NAVY)
        add_text(s, Inches(0.4), y + Inches(0.05), Inches(0.55), Inches(0.35),
                 code, size=14, bold=True, color=MINT, align=PP_ALIGN.CENTER)
        add_text(s, Inches(1.1), y + Inches(0.08), Inches(11.5), Inches(0.4),
                 q, size=13, color=DARK)
    fin(s, "07 Ask")

    # ══════════════════════════════════════════════════════════════════════
    # 26. THREE SENTENCES + SCRIPT
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "08 · 누구에게나 전달할 고정 스크립트", "자문·노트·발표 오프닝에 그대로")
    card(s, Inches(0.4), Inches(1.25), Inches(12.5), Inches(2.4), fill=NAVY)
    add_text(s, Inches(0.65), Inches(1.45), Inches(12), Inches(0.35),
             "세 문장 (복붙용)", size=14, bold=True, color=MINT)
    add_text(s, Inches(0.65), Inches(1.95), Inches(12), Inches(1.5),
             "1. 우리는 IDS 미스센스를 독립 라벨이 아니라 기능 구조를 떠받치는 네트워크의 교란으로 읽는다.\n"
             "2. Ser152 계열은 촉매 자리가 아니라 코어 모듈 사례이며, 개입은 변이 원자가 아니라 인접 표면 패치를 향한다.\n"
             "3. 성공은 결합이 아니라 catalytic competence이며, 현재 구조·계산은 우선순위를 주지만 구제 증명 전이다.",
             size=14, color=WHITE)

    card(s, Inches(0.4), Inches(3.9), Inches(12.5), Inches(2.9), fill=CARD, line=LINE)
    add_text(s, Inches(0.65), Inches(4.05), Inches(12), Inches(0.35),
             "60초 엘리베이터 (화학공학 버전)", size=14, bold=True, color=TEAL)
    add_text(s, Inches(0.65), Inches(4.5), Inches(12), Inches(2.1),
             "헌터 증후군은 리소좀이라는 반응기에서 IDS 효소가 부족해 GAG가 쌓이는 병입니다. 저는 일부 미스센스가\n"
             "‘반응기 칼날 파괴’가 아니라 ‘조립 공장(ER) 수율·품질관리 실패’일 수 있다고 보고, 결정구조 5FQL에 797개 변이를\n"
             "올려 샤페론 적합성 triage 지도를 만들고 있습니다. Ser152처럼 완전히 묻힌 자리는 직접 약이 닿을 수 없어\n"
             "인접 표면을 붙잡는 논리로 설계하고, 성공 지표는 docking이 아니라 기능 구조 회복입니다. 지금은 치료 주장이\n"
             "아니라 의사결정 층을 고정하는 단계이며, ΔΔG와 잔존활성 문헌으로 다음 게이트를 넘으려 합니다.",
             size=13, color=DARK)
    fin(s, "08 Script")

    # ══════════════════════════════════════════════════════════════════════
    # 27. CHEMENG UNIQUE LENS
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "08 · 나만의 사고 렌즈 — 화학공학이 상징적으로 기여하는 지점",
                "이 슬라이드가 ‘남의 프로젝트 요약’이 아니라 ‘내 연구’가 되게 한다")
    lenses = [
        ("물질수지", "기능 IDS 분자 수지\n생성 − 폐기(ERAD) − 미배달\n= 리소좀 재고", "활성 assay를\n‘유량 KPI’로 재해석"),
        ("단위공정", "ER / QC / 골지 / 리소좀을\n직렬 공정으로 모델링", "개입 지점 = 어느 유닛을\n건드릴지 선택"),
        ("열역학", "ΔΔG = 공정 에너지 수지\nAF 그림 ≠ 에너지 상태", "착시를 끊는\n정량 언어"),
        ("Fault mode", "촉매 kill / 코어 불안정 /\n포켓 인접 / PTM 결함", "같은 증상(활성↓)의\n다른 처방전"),
        ("Spec vs 운전", "서열 스펙 vs 접힘 운전창\n(ER pH·샤페론 점유)", "PC = 운전창을\n밀어 주는 보조제"),
        ("Scale-up 윤리", "계산 → 파일럿 → 검증\n게이트 없이 스케일업 금지", "RFdiffusion 대량 =\n파일럿 없는 플랜트"),
    ]
    for i, (t, d, r) in enumerate(lenses):
        col, row = i % 3, i // 3
        left = Inches(0.35 + col * 4.3)
        top = Inches(1.25 + row * 2.85)
        card(s, left, top, Inches(4.1), Inches(2.65), fill=CARD, line=LINE)
        add_rect(s, left, top, Inches(4.1), Inches(0.5), fill=TEAL)
        add_text(s, left + Inches(0.2), top + Inches(0.1), Inches(3.7), Inches(0.35),
                 t, size=14, bold=True, color=WHITE)
        add_text(s, left + Inches(0.2), top + Inches(0.7), Inches(3.7), Inches(1.1),
                 d, size=13, color=DARK)
        add_text(s, left + Inches(0.2), top + Inches(1.9), Inches(3.7), Inches(0.55),
                 r, size=12, bold=True, color=TEAL_DK)
    fin(s, "08 Lens")

    # ══════════════════════════════════════════════════════════════════════
    # 28. FILE MAP
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    section_bar(s, "부록 · 파일 지도 (더 깊게 파고들 때)", "새 도구를 쌓지 말고 한 파일만 더 연다")
    files = [
        ("나침반", "IDSX_한눈에_로직맵_다음스텝.md"),
        ("규칙·claim", "IDS_Chaperone_Amenability_Rules_Claim_Ambroxol_v1.md"),
        ("사고 궤적", "Advisor_Briefing_Package/02_고민의_궤적_논리메모.md"),
        ("질문지", "Advisor_Briefing_Package/03_교수_학자_질문지.md"),
        ("Ser152 전수", "Ser152Ile_Mechanism_Review/00_…전수검토.md"),
        ("AF 보고서", "AF2_runs/AF3_analysis/00_AF3_WT_vs_S152I_분석보고서.md"),
        ("로드맵", "IDS_Chaperone_Binder_RD_Roadmap.md"),
        ("대화 정리", "IDSX_대화정리_안정성_ER_암브록솔_2026-07-11.md"),
        ("AF 한계", "Advisor…/07_AlphaFold_한계점과_패러다임_전환.md"),
        ("공개 포털", "https://thilo-yoo.github.io/idsx-lab/"),
        ("본 패키지", "Desktop/IDSX_Master_Package_2026/"),
        ("원천 폴더", "Desktop/IDSX/"),
    ]
    for i, (k, v) in enumerate(files):
        col, row = i % 2, i // 2
        left = Inches(0.35 + col * 6.5)
        top = Inches(1.2 + row * 0.9)
        card(s, left, top, Inches(6.3), Inches(0.8), fill=CARD, line=LINE)
        add_text(s, left + Inches(0.2), top + Inches(0.2), Inches(1.6), Inches(0.4),
                 k, size=13, bold=True, color=TEAL)
        add_text(s, left + Inches(1.9), top + Inches(0.2), Inches(4.2), Inches(0.4),
                 v, size=12, color=DARK)
    fin(s, "Appendix")

    # ══════════════════════════════════════════════════════════════════════
    # 29. CLOSING
    # ══════════════════════════════════════════════════════════════════════
    s = blank(prs)
    add_rect(s, 0, 0, W, H, fill=NAVY)
    add_text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(0.5),
             "닫는 한 줄", size=16, bold=True, color=MINT)
    add_text(s, Inches(0.6), Inches(2.7), Inches(12), Inches(1.5),
             "치료 성공을 주장하기 전에,\n‘누구에게 어떤 논리로 시도할 가치가 있는가’를\n구조적으로 가르는 지도를 먼저 고정한다.",
             size=24, bold=True, color=WHITE)
    add_text(s, Inches(0.6), Inches(4.8), Inches(12), Inches(1.2),
             "IDSX = Structure-guided decision layer for IDS missense disease\n"
             "P22304 · 5FQL · competence over binding · network over residue worship\n"
             "기록 · 자문 · 협업을 위한 마스터 덱 v2.0 · 2026-07",
             size=14, color=SOFT_TEAL)
    fin(s, "Close")

    # footers
    total = len(built)
    for i, (slide, section) in enumerate(built):
        if i == 0 or i == total - 1:
            continue  # cover & close already full navy
        footer(slide, i + 1, total, section)

    out_pptx = os.path.join(OUT_DIR, "IDSX_Master_Advisor_Presentation.pptx")
    prs.save(out_pptx)
    print(f"Saved: {out_pptx}")
    print(f"Slides: {total}")
    return out_pptx, total


if __name__ == "__main__":
    build()
